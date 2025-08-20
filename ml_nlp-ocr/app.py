import os
import re
import time
import logging
import traceback
import json
import pytesseract  # For OCR
from PIL import Image, ImageEnhance, ImageFilter
import fitz  # PyMuPDF for handling PDFs
from docling.document_converter import DocumentConverter
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import io
import cv2
import numpy as np
import tempfile
import concurrent.futures
import shutil
from datetime import datetime, timedelta
import signal
import sys
from tqdm import tqdm

# Configure logging with rotation to avoid huge log files
import logging.handlers
log_handler = logging.handlers.RotatingFileHandler(
    'file_processing.log', maxBytes=10*1024*1024, backupCount=5)
logging_format = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
log_handler.setFormatter(logging_format)
logger = logging.getLogger()
logger.setLevel(logging.INFO)
logger.addHandler(log_handler)

# Add console handler for immediate feedback
console_handler = logging.StreamHandler()
console_handler.setFormatter(logging_format)
logger.addHandler(console_handler)

# Configuration paths
UPLOAD_FOLDER = "uploaded_files"
PROCESSED_FOLDER = "processed_docs"
IMAGE_FOLDER = "extracted_images"
FAILED_FOLDER = "failed_processing"  # For files that couldn't be processed
TRACKING_FILE = "processed_files.json"  # To track which files have been processed
CONFIG_FILE = "processor_config.json"  # For configuration settings

# Create necessary directories
for folder in [UPLOAD_FOLDER, PROCESSED_FOLDER, IMAGE_FOLDER, FAILED_FOLDER]:
    os.makedirs(folder, exist_ok=True)

# Default configuration
DEFAULT_CONFIG = {
    "scan_interval_minutes": 5,  # Check for missed files every 5 minutes
    "ocr_languages": "eng",      # Default OCR language
    "max_workers": max(1, (os.cpu_count() or 4) - 1),  # Leave one core free
    "monitor_mode": "both",      # 'watch', 'scan', or 'both'
    "image_quality_factor": 1.5, # Enhancement factor for images
    "last_scan_time": None       # Last time a full scan was performed
}

# Configure Tesseract path based on OS
import platform
if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"
# On Linux/Mac, pytesseract will use the system installation

# Global variables
processed_files = {}  # Will be loaded from tracking file
config = DEFAULT_CONFIG.copy()  # Will be loaded/updated from config file
processing_lock = False  # Simple lock to prevent parallel scans

# Signal handler for graceful shutdown
def signal_handler(sig, frame):
    print("\nShutting down gracefully. Saving tracking data...")
    save_processed_files()
    save_config()
    print("Shutdown complete. Exiting.")
    sys.exit(0)

# Register signal handlers
signal.signal(signal.SIGINT, signal_handler)  # Ctrl+C
signal.signal(signal.SIGTERM, signal_handler)  # Termination signal

def load_config():
    """Load configuration from file or create default if not exists."""
    global config
    try:
        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, 'r') as f:
                loaded_config = json.load(f)
                # Update config with loaded values, keeping defaults for missing keys
                for key, value in loaded_config.items():
                    config[key] = value
            logging.info(f"Configuration loaded from {CONFIG_FILE}")
        else:
            save_config()  # Create default config file
            logging.info(f"Created default configuration in {CONFIG_FILE}")
    except Exception as e:
        logging.error(f"Error loading configuration: {e}")

def save_config():
    """Save current configuration to file."""
    try:
        # Update last scan time
        config["last_scan_time"] = datetime.now().isoformat()
        
        with open(CONFIG_FILE, 'w') as f:
            json.dump(config, f, indent=2)
    except Exception as e:
        logging.error(f"Error saving configuration: {e}")

def load_processed_files():
    """Load the list of processed files from tracking file."""
    global processed_files
    try:
        if os.path.exists(TRACKING_FILE):
            with open(TRACKING_FILE, 'r') as f:
                processed_files = json.load(f)
            logging.info(f"Loaded {len(processed_files)} processed file records")
        else:
            processed_files = {}
            save_processed_files()  # Create empty tracking file
    except Exception as e:
        logging.error(f"Error loading processed files tracking: {e}")
        processed_files = {}

def save_processed_files():
    """Save the current list of processed files to tracking file."""
    try:
        with open(TRACKING_FILE, 'w') as f:
            json.dump(processed_files, f, indent=2)
    except Exception as e:
        logging.error(f"Error saving processed files tracking: {e}")

def mark_as_processed(file_path, success=True, metadata=None):
    """Mark a file as processed in the tracking database."""
    abs_path = os.path.abspath(file_path)
    file_name = os.path.basename(file_path)
    
    processed_files[abs_path] = {
        "filename": file_name,
        "processed_time": datetime.now().isoformat(),
        "success": success,
        "metadata": metadata or {}
    }
    
    # Save to tracking file periodically (not on every file to reduce I/O)
    if len(processed_files) % 10 == 0:  
        save_processed_files()

def is_already_processed(file_path):
    """Check if a file has already been successfully processed."""
    abs_path = os.path.abspath(file_path)
    return abs_path in processed_files and processed_files[abs_path].get("success", False)

def check_tesseract_installation():
    """Verify that Tesseract is properly installed and configured."""
    try:
        pytesseract.get_tesseract_version()
        logging.info(f"Tesseract version: {pytesseract.get_tesseract_version()}")
        return True
    except Exception as e:
        logging.error(f"Tesseract not properly configured: {e}")
        print(f"ERROR: Tesseract OCR is not properly configured: {e}")
        print("Please ensure Tesseract is installed and its path is correctly set.")
        return False

def preprocess_image(image):
    """
    Preprocess images for better OCR results with multiple enhancement techniques.
    Returns both the enhanced color image (for saving) and the binary version (for OCR).
    """
    try:
        # Create a copy to avoid modifying the original
        enhanced = image.copy()
        
        # Enhance contrast 
        enhancer = ImageEnhance.Contrast(enhanced)
        enhanced = enhancer.enhance(config["image_quality_factor"])
        
        # Enhance sharpness
        enhancer = ImageEnhance.Sharpness(enhanced)
        enhanced = enhancer.enhance(config["image_quality_factor"])
        
        # Convert to numpy array for OpenCV processing
        image_array = np.array(enhanced)
        
        # If grayscale, convert to RGB first
        if len(image_array.shape) == 2:
            image_array = cv2.cvtColor(image_array, cv2.COLOR_GRAY2RGB)
            
        # Convert to grayscale for binary processing
        gray = cv2.cvtColor(image_array, cv2.COLOR_RGB2GRAY)
        
        # Apply multiple noise reduction techniques
        blurred = cv2.GaussianBlur(gray, (3, 3), 0)
        
        # Enhance text using morphological operations
        kernel = np.ones((1, 1), np.uint8)
        morph = cv2.morphologyEx(blurred, cv2.MORPH_CLOSE, kernel)
        
        # Apply adaptive thresholding - works better than Otsu for documents
        binary = cv2.adaptiveThreshold(
            morph, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
            cv2.THRESH_BINARY, 11, 2
        )
        
        # Apply dilation to fill in broken text
        kernel = np.ones((1, 1), np.uint8)
        binary = cv2.dilate(binary, kernel, iterations=1)
        
        # Create a binary PIL image for OCR
        binary_pil = Image.fromarray(binary)
        
        return enhanced, binary_pil
    
    except Exception as e:
        logging.error(f"Error preprocessing image: {e}")
        # If preprocessing fails, return the original image
        return image, image

def save_extracted_image(image, base_name, index):
    """
    Save the extracted image to IMAGE_FOLDER with a unique name.
    Creates subdirectory based on the base_name for better organization.
    """
    try:
        # Create subfolder for this document
        safe_base = re.sub(r'[^a-zA-Z0-9_\-]', '_', base_name)
        img_subfolder = os.path.join(IMAGE_FOLDER, safe_base)
        os.makedirs(img_subfolder, exist_ok=True)
        
        # Save with a sequential index
        img_filename = f"image_{index:03d}.png"
        img_path = os.path.join(img_subfolder, img_filename)
        
        # Save the image
        image.save(img_path)
        
        # Return relative path for markdown reference
        return os.path.join(img_subfolder, img_filename)
    
    except Exception as e:
        logging.error(f"Error saving image {base_name}_{index}: {e}")
        return None

def extract_text_with_docling(file_path):
    """
    Use Docling to convert a file (pdf/docx/pptx) to Markdown with placeholders.
    This approach preserves document structure and inserts placeholders for images.
    """
    try:
        converter = DocumentConverter()
        result = converter.convert(file_path)
        markdown = result.document.export_to_markdown()
        
        # Convert any [IMAGE] markers to <IMAGE_PLACEHOLDER> for consistency
        markdown = re.sub(r'\[IMAGE\]', '<IMAGE_PLACEHOLDER>', markdown)
        
        return markdown
    except Exception as e:
        logging.error(f"Docling conversion error for {file_path}: {e}")
        # Fall back to basic text extraction if Docling fails
        return f"*Error converting document with Docling. Basic text extraction follows:*\n\n"

def perform_ocr(image, lang=None, config='--psm 6'):
    """
    Perform OCR on the given image with error handling.
    Returns empty string if OCR fails.
    """
    try:
        # Use configured languages or default
        if lang is None:
            lang = config.get("ocr_languages", "eng")
            
        # Try to improve OCR with image preprocessing
        _, binary_image = preprocess_image(image)
        
        # Perform OCR with confidence
        ocr_data = pytesseract.image_to_data(
            binary_image, lang=lang, config=config, output_type=pytesseract.Output.DICT
        )
        
        # Filter out low confidence results
        text_parts = []
        for i, conf in enumerate(ocr_data['conf']):
            if float(conf) > 30:  # Only include text with confidence above 30%
                text = ocr_data['text'][i]
                if text.strip():
                    text_parts.append(text)
        
        # Join text parts
        text = " ".join(text_parts)
        
        # Clean up text
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    except Exception as e:
        logging.error(f"OCR error: {e}")
        return ""

def replace_placeholder(markdown_text, image_path, ocr_text):
    """
    Replace the first occurrence of <IMAGE_PLACEHOLDER> with
    the actual markdown reference and OCR text.
    If no placeholder is found, append at the end.
    """
    # Clean OCR text to avoid empty results
    ocr_text = ocr_text.strip() if ocr_text else "*No text detected in image*"
    
    # Create the markdown snippet with image and OCR text
    snippet = f"![Image]({image_path})\n\n*Image text:* {ocr_text}\n\n"
    
    # Replace placeholder if found
    if "<IMAGE_PLACEHOLDER>" in markdown_text:
        return markdown_text.replace("<IMAGE_PLACEHOLDER>", snippet, 1)
    else:
        return markdown_text + "\n\n" + snippet

def extract_text_from_pdf(file_path):
    """
    Extract text and images from PDF using PyMuPDF with Docling structure.
    Uses OCR for images and replaces placeholders in the document.
    """
    base_name = os.path.basename(file_path).rsplit('.', 1)[0]
    
    # First get the structured markdown from Docling
    markdown_text = extract_text_with_docling(file_path)
    
    try:
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(file_path)
        img_counter = 0
        
        # Process each page
        for page_index in range(len(pdf_document)):
            page = pdf_document.load_page(page_index)
            
            # Extract images
            image_list = page.get_images(full=True)
            
            # Process each image
            for img_info in image_list:
                try:
                    xref = img_info[0]
                    base_image = fitz.Pixmap(pdf_document, xref)
                    
                    # Convert to PIL Image
                    if base_image.n < 5:  # RGB
                        pil_img = Image.frombytes("RGB", [base_image.width, base_image.height], base_image.samples)
                    else:  # CMYK
                        cmyk_image = fitz.Pixmap(fitz.csRGB, base_image)
                        pil_img = Image.frombytes("RGB", [cmyk_image.width, cmyk_image.height], cmyk_image.samples)
                        cmyk_image = None  # Free memory
                    
                    # Skip small images (likely icons or decorations)
                    if pil_img.width < 50 or pil_img.height < 50:
                        continue
                    
                    # Enhance and save the image
                    enhanced_img, _ = preprocess_image(pil_img)
                    image_path = save_extracted_image(enhanced_img, base_name, img_counter)
                    
                    # Extract text from image using OCR
                    ocr_text = perform_ocr(enhanced_img)
                    
                    # Replace placeholder with image reference and OCR text
                    markdown_text = replace_placeholder(markdown_text, image_path, ocr_text)
                    img_counter += 1
                    
                except Exception as e:
                    logging.error(f"Error processing image {img_counter} in {file_path}: {e}")
                    continue
            
            # Check if page has text but few or no extracted images
            # This might indicate image-based PDFs that need page-level OCR
            if len(image_list) < 2 and len(page.get_text().strip()) < 100:
                try:
                    # Render page as image
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    
                    # Save the page image
                    image_path = save_extracted_image(img, base_name, f"page_{page_index}")
                    
                    # Perform OCR on the page
                    ocr_text = perform_ocr(img)
                    
                    # Add page content with image and OCR text
                    page_note = f"\n\n## Page {page_index + 1}\n\n"
                    page_note += f"![Page {page_index + 1}]({image_path})\n\n"
                    page_note += f"*Page text:* {ocr_text}\n\n"
                    
                    # Add to markdown
                    markdown_text += page_note
                    
                except Exception as e:
                    logging.error(f"Error processing page {page_index} in {file_path}: {e}")
        
        pdf_document.close()
        return markdown_text
        
    except Exception as e:
        logging.error(f"Error extracting from PDF {file_path}: {str(e)}\n{traceback.format_exc()}")
        return markdown_text + f"\n\n*Error extracting content: {str(e)}*"

def extract_text_from_docx(file_path):
    """
    Extract text and images from DOCX using Docling and python-docx.
    """
    from docx import Document
    base_name = os.path.basename(file_path).rsplit('.', 1)[0]
    
    # Get structured markdown from Docling
    markdown_text = extract_text_with_docling(file_path)
    
    try:
        # Open document to extract images
        doc = Document(file_path)
        img_counter = 0
        
        # Process all relationships to find images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    # Get image binary data
                    image_data = rel.target_part.blob
                    
                    # Save to temp file
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_file:
                        temp_file.write(image_data)
                        tmp_path = temp_file.name
                    
                    try:
                        # Open and process image
                        img = Image.open(tmp_path)
                        
                        # Skip small images (likely icons or decorations)
                        if img.width < 50 or img.height < 50:
                            os.remove(tmp_path)
                            continue
                            
                        # Enhance and save image
                        enhanced_img, _ = preprocess_image(img)
                        image_path = save_extracted_image(enhanced_img, base_name, img_counter)
                        
                        # Extract text from image using OCR
                        ocr_text = perform_ocr(enhanced_img)
                        
                        # Replace placeholder with image reference and OCR text
                        markdown_text = replace_placeholder(markdown_text, image_path, ocr_text)
                        img_counter += 1
                        
                    finally:
                        # Clean up temporary file
                        if os.path.exists(tmp_path):
                            os.remove(tmp_path)
                
                except Exception as e:
                    logging.error(f"Error processing image in {file_path}: {e}")
                    continue
        
        return markdown_text
        
    except Exception as e:
        logging.error(f"Error extracting from DOCX {file_path}: {str(e)}\n{traceback.format_exc()}")
        return markdown_text + f"\n\n*Error extracting content: {str(e)}*"

def extract_text_from_pptx(file_path):
    """
    Extract text and images from PPTX using Docling and python-pptx.
    """
    from pptx import Presentation
    base_name = os.path.basename(file_path).rsplit('.', 1)[0]
    
    # Get structured markdown from Docling
    markdown_text = extract_text_with_docling(file_path)
    
    try:
        # Open presentation
        prs = Presentation(file_path)
        img_counter = 0
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides):
            # Process each shape in the slide
            for shape in slide.shapes:
                # Check if shape is a picture (13) or has a picture fill (1)
                if shape.shape_type in [1, 13] and hasattr(shape, 'image'):
                    try:
                        # Extract image data
                        image = shape.image
                        
                        # Save to temp file
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_f:
                            temp_f.write(image.blob)
                            tmp_path = temp_f.name
                        
                        try:
                            # Open and process image
                            img = Image.open(tmp_path)
                            
                            # Skip small images (likely icons or decorations)
                            if img.width < 50 or img.height < 50:
                                os.remove(tmp_path)
                                continue
                                
                            # Enhance and save image
                            enhanced_img, _ = preprocess_image(img)
                            image_path = save_extracted_image(enhanced_img, base_name, img_counter)
                            
                            # Extract text from image using OCR
                            ocr_text = perform_ocr(enhanced_img)
                            
                            # Replace placeholder with image reference and OCR text
                            markdown_text = replace_placeholder(markdown_text, image_path, ocr_text)
                            img_counter += 1
                            
                        finally:
                            # Clean up temporary file
                            if os.path.exists(tmp_path):
                                os.remove(tmp_path)
                    
                    except Exception as e:
                        logging.error(f"Error processing image in slide {slide_num} of {file_path}: {e}")
                        continue
        
        return markdown_text
        
    except Exception as e:
        logging.error(f"Error extracting from PPTX {file_path}: {str(e)}\n{traceback.format_exc()}")
        return markdown_text + f"\n\n*Error extracting content: {str(e)}*"

def extract_text_from_txt(file_path):
    """Extract text from plain text files with encoding detection."""
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                    # Successfully read, return text
                    return text
            except UnicodeDecodeError:
                # Try next encoding
                continue
        
        # If all encodings fail, try binary mode and decode with error handling
        with open(file_path, "rb") as f:
            binary = f.read()
            text = binary.decode('utf-8', errors='replace')
            return text
            
    except Exception as e:
        logging.error(f"Error reading text file {file_path}: {e}")
        return f"*Error reading text file: {str(e)}*"

def extract_text_from_csv(file_path):
    """Extract text from CSV files with intelligent formatting."""
    import csv
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                markdown = "| "
                row_count = 0
                
                with open(file_path, "r", encoding=encoding) as csv_file:
                    reader = csv.reader(csv_file)
                    headers = next(reader, None)
                    
                    # If we have headers, add them as the first row
                    if headers:
                        markdown += " | ".join(headers) + " |\n| "
                        markdown += " | ".join(["---"] * len(headers)) + " |\n"
                    
                    # Add data rows
                    for row in reader:
                        row_count += 1
                        if row_count > 100:  # Limit to 100 rows to avoid huge files
                            markdown += "\n\n*...truncated, showing first 100 rows...*"
                            break
                        markdown += " | ".join(row) + " |\n"
                
                # Successfully read, return markdown table
                return markdown
                
            except UnicodeDecodeError:
                # Try next encoding
                continue
        
        # If all encodings fail, use a basic approach
        with open(file_path, "rb") as f:
            binary = f.read()
            text = binary.decode('utf-8', errors='replace')
            return f"```\n{text}\n```"
            
    except Exception as e:
        logging.error(f"Error reading CSV file {file_path}: {e}")
        return f"*Error reading CSV file: {str(e)}*"

def extract_text_from_image(file_path):
    """Extract text from image files directly."""
    try:
        # Open the image
        img = Image.open(file_path)
        
        # Save a copy to the image folder
        base_name = os.path.basename(file_path).rsplit('.', 1)[0]
        image_path = save_extracted_image(img, base_name, 0)
        
        # Process image for OCR
        enhanced_img, binary_img = preprocess_image(img)
        
        # Perform OCR
        ocr_text = perform_ocr(binary_img)
        
        # Create markdown with image and text
        markdown = f"# Image: {base_name}\n\n"
        markdown += f"![Image]({image_path})\n\n"
        markdown += f"## Extracted Text\n\n{ocr_text}\n"
        
        return markdown
        
    except Exception as e:
        logging.error(f"Error processing image file {file_path}: {e}")
        return f"*Error processing image file: {str(e)}*"

def process_file(file_path):
    """
    Process a single document file by extracting text and images.
    Returns True if processing was successful, False otherwise.
    """
    if not os.path.exists(file_path):
        logging.error(f"File does not exist: {file_path}")
        return False
        
    # Skip if already processed
    if is_already_processed(file_path):
        logging.info(f"Skipping already processed file: {file_path}")
        return True
        
    file_name = os.path.basename(file_path)
    extension = file_name.split('.')[-1].lower()
    
    logging.info(f"Processing started: {file_name}")
    print(f"Processing: {file_name}")
    
    try:
        # Choose the appropriate extraction method based on file type
        if extension == 'pdf':
            text_md = extract_text_from_pdf(file_path)
        elif extension == 'docx':
            text_md = extract_text_from_docx(file_path)
        elif extension in ['pptx', 'ppt']:
            text_md = extract_text_from_pptx(file_path)
        elif extension == 'txt':
            text_md = extract_text_from_txt(file_path)
        elif extension == 'csv':
            text_md = extract_text_from_csv(file_path)
        elif extension in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif']:
            text_md = extract_text_from_image(file_path)
        else:
            # Try using Docling for unsupported formats
            logging.warning(f"Attempting to process unsupported format: {file_name}")
            try:
                text_md = extract_text_with_docling(file_path)
                if not text_md:
                    raise ValueError("Docling returned empty result")
            except Exception as e:
                logging.error(f"Cannot process unsupported file: {file_name}: {e}")
                # Move to failed folder
                failed_path = os.path.join(FAILED_FOLDER, file_name)
                shutil.copy2(file_path, failed_path)
                mark_as_processed(file_path, success=False, metadata={"error": str(e)})
                return False

        # Add metadata and save to markdown
        if save_to_markdown(file_name, text_md):
            mark_as_processed(file_path, success=True, 
                              metadata={"file_type": extension, "characters": len(text_md)})
            logging.info(f"Processing completed successfully: {file_name}")
            return True
        else:
            mark_as_processed(file_path, success=False, 
                              metadata={"error": "Failed to save markdown"})
            return False

    except Exception as e:
        logging.error(f"Error processing file {file_name}: {str(e)}\n{traceback.format_exc()}")
        # Move to failed folder
        failed_path = os.path.join(FAILED_FOLDER, file_name)
        shutil.copy2(file_path, failed_path)
        mark_as_processed(file_path, success=False, metadata={"error": str(e)})
        return False

def save_to_markdown(file_name, text):
    """
    Saves extracted text to a .md file with sanitized name and metadata.
    Returns True if successful, False otherwise.
    """
    try:
        # Create safe filename
        safe_name = re.sub(r'[^a-zA-Z0-9_\-]', '_', file_name)
        md_path = os.path.join(PROCESSED_FOLDER, f"{safe_name}.md")
        
        # Add metadata header
        metadata = f"""---
title: "{file_name}"
source_file: "{file_name}"
extraction_date: "{time.strftime('%Y-%m-%d %H:%M:%S')}"
plagiarism_check_status: "pending"
---

# {file_name}

"""
        
        # Combine metadata with content
        full_content = metadata + text
        
        # Save to file
        with open(md_path, 'w', encoding='utf-8') as md_file:
            md_file.write(full_content)
            
        print(f"✅ Processed and saved: {md_path}")
        return True
        
    except Exception as e:
        logging.error(f"Error saving markdown file for {file_name}: {str(e)}")
        return False

def scan_for_new_files():
    """
    Scan upload folder for new files not yet processed.
    Returns the number of files processed.
    """
    global processing_lock
    
    # Skip if already processing
    if processing_lock:
        logging.info("Skipping scan as another scan is in progress")
        return 0
        
    try:
        processing_lock = True
        
        # Get list of all files in upload folder
        all_files = [os.path.join(UPLOAD_FOLDER, f) for f in os.listdir(UPLOAD_FOLDER) 
                    if os.path.isfile(os.path.join(UPLOAD_FOLDER, f))]
        
        # Filter to only new files
        new_files = [f for f in all_files if not is_already_processed(f)]
        
        if not new_files:
            logging.info("No new files found to process")
            processing_lock = False
            return 0
            
        logging.info(f"Found {len(new_files)} new files to process")
        
        # Process files in parallel
        processed_count = 0
        with concurrent.futures.ThreadPoolExecutor(max_workers=config["max_workers"]) as executor:
            # Use tqdm for progress tracking
            with tqdm(total=len(new_files), desc="Processing Files", unit="file") as progress:
                # Submit all tasks
                future_to_file = {executor.submit(process_file, file_path): file_path for file_path in new_files}
                
                # Process as they complete
                for future in concurrent.futures.as_completed(future_to_file):
                    file_path = future_to_file[future]
                    try:
                        if future.result():
                            processed_count += 1
                    except Exception as e:
                        logging.error(f"Executor error for {file_path}: {e}")
                    finally:
                        progress.update(1)
        
        # Save tracking data
        save_processed_files()
        
        logging.info(f"Scan completed. Processed {processed_count} of {len(new_files)} files.")
        processing_lock = False
        return processed_count
        
    except Exception as e:
        logging.error(f"Error during scan: {str(e)}\n{traceback.format_exc()}")
        processing_lock = False
        return 0

class FileHandler(FileSystemEventHandler):
    """
    Watchdog handler for real-time file monitoring.
    """
    def on_created(self, event):
        """Handle file creation events."""
        if event.is_directory:
            return
            
        file_path = event.src_path
        
        # Check if this is a file in our upload directory
        if os.path.dirname(os.path.abspath(file_path)) == os.path.abspath(UPLOAD_FOLDER):
            logging.info(f"New file detected: {os.path.basename(file_path)}")
            
            # Wait a moment for file to be completely written
            time.sleep(2)
            
            # Process the file
            try:
                process_file(file_path)
                # Save tracking data after each file in watch mode
                save_processed_files()
            except Exception as e:
                logging.error(f"Error processing new file {file_path}: {e}")

def start_file_watcher():
    """
    Start the watchdog observer to monitor the upload folder.
    """
    logging.info(f"Starting file watcher on folder: {UPLOAD_FOLDER}")
    event_handler = FileHandler()
    observer = Observer()
    observer.schedule(event_handler, path=UPLOAD_FOLDER, recursive=False)
    observer.start()
    return observer

def run_periodic_scan():
    """
    Periodically scan for missed files.
    This handles files that might have been added while the program was not running.
    """
    while True:
        logging.info("Running periodic scan for missed files")
        scan_for_new_files()
        
        # Update last scan time in config
        config["last_scan_time"] = datetime.now().isoformat()
        save_config()
        
        # Sleep for the configured interval
        interval_minutes = config.get("scan_interval_minutes", 5)
        time.sleep(interval_minutes * 60)

def main():
    """Main execution function."""
    print("Document Processing System Starting...")
    print(f"Upload folder: {os.path.abspath(UPLOAD_FOLDER)}")
    print(f"Processed folder: {os.path.abspath(PROCESSED_FOLDER)}")
    
    # Verify Tesseract installation
    if not check_tesseract_installation():
        print("WARNING: Proceeding without proper Tesseract OCR configuration.")
        print("Images may not be properly processed.")
    
    # Load configuration and tracking data
    load_config()
    load_processed_files()
    
    # Determine monitoring mode
    mode = config.get("monitor_mode", "both")
    
    try:
        # Start file watcher if in 'watch' or 'both' mode
        observer = None
        if mode in ["watch", "both"]:
            observer = start_file_watcher()
            print("File watcher started - monitoring for new files in real-time")
        
        # Initial scan for any files that might have been missed
        scan_count = scan_for_new_files()
        print(f"Initial scan completed: {scan_count} files processed")
        
        # Start periodic scan thread if in 'scan' or 'both' mode
        scan_thread = None
        if mode in ["scan", "both"]:
            scan_thread = concurrent.futures.ThreadPoolExecutor(max_workers=1)
            scan_thread.submit(run_periodic_scan)
            print(f"Periodic scan enabled - checking every {config['scan_interval_minutes']} minutes")
            
        print("\nDocument Processor is running. Press Ctrl+C to exit.")
        
        # Keep the main thread running
        try:
            while True:
                time.sleep(1)
        except KeyboardInterrupt:
            print("\nShutdown requested...")
        
        # Clean shutdown
        if observer:
            observer.stop()
            observer.join()
        
        if scan_thread:
            scan_thread.shutdown(wait=False)
            
        save_processed_files()
        save_config()
        print("Document Processor shut down successfully.")
        
    except Exception as e:
        logging.error(f"Critical error in main execution: {str(e)}\n{traceback.format_exc()}")
        print(f"Critical error: {str(e)}")
        print("Check log file for details: file_processing.log")
        
        # Save data before exit
        save_processed_files()
        save_config()

if __name__ == "__main__":
    main()


# in the classroom ananlyse the code files and structure and everything and help me make changes so the
# files when stored are noted the order in ehich they are uploaded so that if later files if uploaded are
# same that is plagiarised then it can be said it is similar to previous submissions for such matching use
# the matching algorithm in the Matcher_algo\plag-detect.py and check if they are similar and also refer
# ml_nlp-ocr\app.py for handling files when uploaded in the website. also make it modular as later we will
# also integrate another model to t=detect if the uploaded assignments are ai-geenrated, so steps -
# upload handled, upload files preprocessed and stored and logged and ordered and anything else we might
# need, then uploaded files compared using the plag-detecct algo by sending data thought api as it will
# also need to send data to another module for ai detection which will be done later