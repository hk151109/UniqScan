from flask import Flask, request, jsonify
import os
import json
import logging
import nltk
import time
import random
import hashlib
import requests
import tempfile
import shutil
from datetime import datetime
from difflib import SequenceMatcher
from nltk.metrics.distance import edit_distance as editDistance
from nltk.stem.lancaster import LancasterStemmer
from nltk.util import ngrams
from pathlib import Path
import re
from urllib.parse import urlparse
import mimetypes
from werkzeug.utils import secure_filename

# PDF/OCR Processing imports
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import fitz  # PyMuPDF for handling PDFs
import cv2
import numpy as np
import io
import traceback
from docx import Document
from pptx import Presentation
import csv

# =====================================================================
# CONFIGURABLE PARAMETERS - Easily adjust these values
# =====================================================================
THRESHOLD = 3     # Minimum match size for initial matching
CUTOFF = 5        # Minimum match size after extending
NGRAM = 3         # N-gram size for tokenization
DISTANCE = 8      # Minimum distance between matches for healing

# File directories - where to store downloaded and processed files
DOWNLOAD_DIR = "downloads"    # Downloaded files
TEXT_DIR = "extracted_text"   # Extracted text files
REPORTS_DIR = "reports"       # Similarity reports
LOGS_DIR = "logs"             # Log files

# Create necessary directories
for directory in [DOWNLOAD_DIR, TEXT_DIR, REPORTS_DIR, LOGS_DIR]:
    os.makedirs(directory, exist_ok=True)

# External AI score API endpoint (can be overridden via environment variable)
AI_SCORE_API_URL = f"{os.getenv('AI_SCORE_API_URL', 'http://localhost:8001').rstrip('/')}/classify"

# Initialize Flask app
app = Flask(__name__)

# Configure logging
log_file = os.path.join(LOGS_DIR, f"plagiarism_detection_{datetime.now().strftime('%Y%m%d')}.log")
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)

# Global storage for user scores
user_scores_log = {}

# Configure Tesseract path based on OS
import platform
if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"

class PDFProcessor:
    """Enhanced PDF processing with OCR capabilities"""
    
    def __init__(self):
        print("Initializing PDF processor...")
        self.check_tesseract_installation()
    
    def check_tesseract_installation(self):
        """Verify that Tesseract is properly installed and configured."""
        try:
            pytesseract.get_tesseract_version()
            logging.info(f"Tesseract version: {pytesseract.get_tesseract_version()}")
            print(f"✓ Tesseract OCR ready: {pytesseract.get_tesseract_version()}")
            return True
        except Exception as e:
            logging.warning(f"Tesseract not properly configured: {e}")
            logging.warning("OCR functionality may be limited")
            print(f"⚠ Tesseract warning: {e}")
            print("  OCR functionality may be limited")
            return False
    
    def preprocess_image(self, image, quality_factor=1.5):
        """
        Preprocess images for better OCR results with multiple enhancement techniques.
        Returns both the enhanced color image and the binary version for OCR.
        """
        try:
            print("Preprocessing image for OCR...")
            # Create a copy to avoid modifying the original
            enhanced = image.copy()
            
            # Enhance contrast 
            enhancer = ImageEnhance.Contrast(enhanced)
            enhanced = enhancer.enhance(quality_factor)
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(enhanced)
            enhanced = enhancer.enhance(quality_factor)
            
            # Convert to numpy array for OpenCV processing
            image_array = np.array(enhanced)
            
            # If grayscale, convert to RGB first
            if len(image_array.shape) == 2:
                image_array = cv2.cvtColor(image_array, cv2.COLOR_GRAY2RGB)
                
            # Convert to grayscale for binary processing
            gray = cv2.cvtColor(image_array, cv2.COLOR_RGB2GRAY)
            
            # Apply multiple noise reduction techniques
            blurred = cv2.GaussianBlur(gray, (3, 3), 0)
            
            # Enhance text using morphological operations
            kernel = np.ones((1, 1), np.uint8)
            morph = cv2.morphologyEx(blurred, cv2.MORPH_CLOSE, kernel)
            
            # Apply adaptive thresholding - works better than Otsu for documents
            binary = cv2.adaptiveThreshold(
                morph, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                cv2.THRESH_BINARY, 11, 2
            )
            
            # Apply dilation to fill in broken text
            kernel = np.ones((1, 1), np.uint8)
            binary = cv2.dilate(binary, kernel, iterations=1)
            
            # Create a binary PIL image for OCR
            binary_pil = Image.fromarray(binary)
            
            print("Image preprocessing complete.")
            return enhanced, binary_pil
        
        except Exception as e:
            logging.error(f"Error preprocessing image: {e}")
            print(f"Error preprocessing image: {e}")
            # If preprocessing fails, return the original image
            return image, image
    
    def perform_ocr(self, image, lang="eng", config='--psm 6'):
        """
        Perform OCR on the given image with error handling.
        Returns empty string if OCR fails.
        """
        try:
            print("Performing OCR on image...")
            # Try to improve OCR with image preprocessing
            _, binary_image = self.preprocess_image(image)
            
            # Perform OCR with confidence
            ocr_data = pytesseract.image_to_data(
                binary_image, lang=lang, config=config, output_type=pytesseract.Output.DICT
            )
            
            # Filter out low confidence results
            text_parts = []
            for i, conf in enumerate(ocr_data['conf']):
                if float(conf) > 30:  # Only include text with confidence above 30%
                    text = ocr_data['text'][i]
                    if text.strip():
                        text_parts.append(text)
            
            # Join text parts
            text = " ".join(text_parts)
            
            # Clean up text
            text = re.sub(r'\s+', ' ', text).strip()
            
            print(f"OCR complete. Extracted {len(text)} characters.")
            return text
        except Exception as e:
            logging.error(f"OCR error: {e}")
            print(f"OCR error: {e}")
            return ""
    
    def extract_text_from_pdf(self, file_path):
        """
        Extract text and images from PDF using PyMuPDF with OCR for images.
        """
        try:
            print(f"Extracting text from PDF: {os.path.basename(file_path)}")
            # Open PDF with PyMuPDF
            pdf_document = fitz.open(file_path)
            extracted_text = ""
            
            # Process each page
            for page_index in range(len(pdf_document)):
                print(f"Processing page {page_index + 1} of {len(pdf_document)}...")
                page = pdf_document.load_page(page_index)
                
                # Extract direct text first
                page_text = page.get_text()
                
                # If page has substantial text, use it
                if len(page_text.strip()) > 50:
                    extracted_text += f"\n\n--- Page {page_index + 1} ---\n\n"
                    extracted_text += page_text
                else:
                    # Page might be image-based, try OCR on the entire page
                    try:
                        print(f"Page {page_index + 1} has limited text, trying OCR...")
                        # Render page as image
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        
                        # Perform OCR on the page
                        ocr_text = self.perform_ocr(img)
                        
                        if ocr_text.strip():
                            extracted_text += f"\n\n--- Page {page_index + 1} (OCR) ---\n\n"
                            extracted_text += ocr_text
                        
                    except Exception as e:
                        logging.error(f"Error processing page {page_index} with OCR: {e}")
                        print(f"Error processing page {page_index} with OCR: {e}")
                
                # Extract images and perform OCR on them
                image_list = page.get_images(full=True)
                
                if image_list:
                    print(f"Found {len(image_list)} images on page {page_index + 1}.")
                
                for img_index, img_info in enumerate(image_list):
                    try:
                        print(f"Processing image {img_index + 1} on page {page_index + 1}...")
                        xref = img_info[0]
                        base_image = fitz.Pixmap(pdf_document, xref)
                        
                        # Convert to PIL Image
                        if base_image.n < 5:  # RGB
                            pil_img = Image.frombytes("RGB", [base_image.width, base_image.height], base_image.samples)
                        else:  # CMYK
                            cmyk_image = fitz.Pixmap(fitz.csRGB, base_image)
                            pil_img = Image.frombytes("RGB", [cmyk_image.width, cmyk_image.height], cmyk_image.samples)
                            cmyk_image = None  # Free memory
                        
                        # Skip small images (likely icons or decorations)
                        if pil_img.width < 50 or pil_img.height < 50:
                            print(f"Skipping small image: {pil_img.width}x{pil_img.height}")
                            continue
                        
                        # Extract text from image using OCR
                        ocr_text = self.perform_ocr(pil_img)
                        
                        if ocr_text.strip():
                            extracted_text += f"\n\n--- Page {page_index + 1}, Image {img_index + 1} (OCR) ---\n\n"
                            extracted_text += ocr_text
                        
                    except Exception as e:
                        logging.error(f"Error processing image {img_index} on page {page_index}: {e}")
                        print(f"Error processing image {img_index} on page {page_index}: {e}")
                        continue
            
            pdf_document.close()
            print(f"PDF extraction complete. Extracted {len(extracted_text)} characters.")
            return extracted_text.strip()
            
        except Exception as e:
            logging.error(f"Error extracting from PDF {file_path}: {str(e)}\n{traceback.format_exc()}")
            print(f"Error extracting from PDF {file_path}: {str(e)}")
            return f"Error extracting content from PDF: {str(e)}"
    
    def extract_text_from_docx(self, file_path):
        """Extract text from DOCX files."""
        try:
            print(f"Extracting text from DOCX: {os.path.basename(file_path)}")
            doc = Document(file_path)
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            result = "\n\n".join(text_content)
            print(f"DOCX extraction complete. Extracted {len(result)} characters.")
            return result
            
        except Exception as e:
            logging.error(f"Error extracting from DOCX {file_path}: {e}")
            print(f"Error extracting from DOCX {file_path}: {e}")
            return f"Error extracting content from DOCX: {str(e)}"
    
    def extract_text_from_pptx(self, file_path):
        """Extract text from PPTX files."""
        try:
            print(f"Extracting text from PPTX: {os.path.basename(file_path)}")
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {slide_num + 1} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            result = "\n\n".join(text_content)
            print(f"PPTX extraction complete. Extracted {len(result)} characters.")
            return result
            
        except Exception as e:
            logging.error(f"Error extracting from PPTX {file_path}: {e}")
            print(f"Error extracting from PPTX {file_path}: {e}")
            return f"Error extracting content from PPTX: {str(e)}"
    
    def extract_text_from_txt(self, file_path):
        """Extract text from plain text files with encoding detection."""
        try:
            print(f"Extracting text from TXT: {os.path.basename(file_path)}")
            # Try multiple encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read()
                        print(f"TXT extraction complete. Extracted {len(content)} characters.")
                        return content
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try binary mode with error handling
            with open(file_path, "rb") as f:
                binary = f.read()
                content = binary.decode('utf-8', errors='replace')
                print(f"TXT extraction complete (using binary mode). Extracted {len(content)} characters.")
                return content
                
        except Exception as e:
            logging.error(f"Error reading text file {file_path}: {e}")
            print(f"Error reading text file {file_path}: {e}")
            return f"Error reading text file: {str(e)}"
    
    def extract_text_from_csv(self, file_path):
        """Extract text from CSV files."""
        try:
            print(f"Extracting text from CSV: {os.path.basename(file_path)}")
            text_content = []
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as csv_file:
                        reader = csv.reader(csv_file)
                        row_count = 0
                        
                        for row in reader:
                            row_count += 1
                            if row_count > 100:  # Limit to 100 rows
                                text_content.append("...truncated, showing first 100 rows...")
                                break
                            text_content.append(" ".join(row))
                    
                    result = "\n".join(text_content)
                    print(f"CSV extraction complete. Extracted {len(result)} characters.")
                    return result
                    
                except UnicodeDecodeError:
                    continue
            
            return "Error: Could not decode CSV file"
            
        except Exception as e:
            logging.error(f"Error reading CSV file {file_path}: {e}")
            print(f"Error reading CSV file {file_path}: {e}")
            return f"Error reading CSV file: {str(e)}"
    
    def extract_text_from_image(self, file_path):
        """Extract text from image files using OCR."""
        try:
            print(f"Extracting text from image: {os.path.basename(file_path)}")
            img = Image.open(file_path)
            ocr_text = self.perform_ocr(img)
            result = ocr_text if ocr_text.strip() else "No text detected in image"
            print(f"Image OCR complete. Extracted {len(result)} characters.")
            return result
            
        except Exception as e:
            logging.error(f"Error processing image file {file_path}: {e}")
            print(f"Error processing image file {file_path}: {e}")
            return f"Error processing image file: {str(e)}"
    
    def process_file_content(self, file_path, content_type=None):
        """
        Process a file and extract text content based on file type.
        """
        try:
            print(f"Processing file: {os.path.basename(file_path)}")
            # Determine file type from extension if content_type not provided
            if not content_type:
                _, ext = os.path.splitext(file_path)
                ext = ext.lower().lstrip('.')
            else:
                # Extract extension from content type or file path
                if '/' in content_type:
                    ext = content_type.split('/')[-1]
                else:
                    _, ext = os.path.splitext(file_path)
                    ext = ext.lower().lstrip('.')
            
            print(f"Detected file type: {ext}")
            
            # Route to appropriate extraction method
            if ext == 'pdf':
                return self.extract_text_from_pdf(file_path)
            elif ext in ['docx', 'doc']:
                return self.extract_text_from_docx(file_path)
            elif ext in ['pptx', 'ppt']:
                return self.extract_text_from_pptx(file_path)
            elif ext == 'txt':
                return self.extract_text_from_txt(file_path)
            elif ext == 'csv':
                return self.extract_text_from_csv(file_path)
            elif ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif']:
                return self.extract_text_from_image(file_path)
            else:
                # Try to read as text file for unknown types
                logging.warning(f"Unknown file type: {ext}. Attempting text extraction.")
                print(f"Unknown file type: {ext}. Attempting text extraction.")
                return self.extract_text_from_txt(file_path)
                
        except Exception as e:
            logging.error(f"Error processing file {file_path}: {e}")
            print(f"Error processing file {file_path}: {e}")
            return f"Error processing file: {str(e)}"


class Text:
    def __init__(self, raw_text, label, filepath, removeStopwords=True):
        print(f"Creating Text object for: {label}")
        if isinstance(raw_text, list):
            self.text = ' \n '.join(raw_text)
        else:
            self.text = raw_text
        self.label = label
        self.filepath = filepath
        self.preprocess(self.text)
        self.tokens = self.getTokens(removeStopwords)
        self.trigrams = self.ngrams(NGRAM)  # Using configurable NGRAM value
        self.checksum = self.calculate_checksum()
        print(f"Text object created: {len(self.tokens)} tokens, {len(self.trigrams)} {NGRAM}-grams")

    def calculate_checksum(self):
        """Calculate a checksum for the file content to detect changes"""
        return hashlib.md5(self.text.encode()).hexdigest()

    def preprocess(self, text):
        """ Heals hyphenated words, and maybe other things. """
        self.text = re.sub(r'([A-Za-z])- ([a-z])', r'\1\2', text)

    def getTokens(self, removeStopwords=True):
        """ Tokenizes the text, breaking it up into words, removing punctuation. """
        tokenizer = nltk.RegexpTokenizer('[a-zA-Z]\\w+\'?\\w*')
        spans = list(tokenizer.span_tokenize(self.text))
        
        if spans:
            self.length = spans[-1][-1]
        else:
            self.length = 0
            
        tokens = tokenizer.tokenize(self.text)
        tokens = [token.lower() for token in tokens]
        stemmer = LancasterStemmer()
        tokens = [stemmer.stem(token) for token in tokens]
        
        if not removeStopwords:
            self.spans = spans
            return tokens
            
        tokenSpans = list(zip(tokens, spans))
        stopwords = nltk.corpus.stopwords.words('english')
        tokenSpans = [token for token in tokenSpans if token[0] not in stopwords]
        self.spans = [x[1] for x in tokenSpans]
        return [x[0] for x in tokenSpans]

    def ngrams(self, n):
        """ Returns ngrams for the text."""
        return list(ngrams(self.tokens, n))


class ExtendedMatch:
    def __init__(self, a, b, sizeA, sizeB):
        self.a = a
        self.b = b
        self.sizeA = sizeA
        self.sizeB = sizeB
        self.healed = False
        self.extendedBackwards = 0
        self.extendedForwards = 0

    def __repr__(self):
        out = "a: %s, b: %s, size a: %s, size b: %s" % (self.a, self.b, self.sizeA, self.sizeB)
        if self.extendedBackwards:
            out += ", extended backwards x%s" % self.extendedBackwards
        if self.extendedForwards:
            out += ", extended forwards x%s" % self.extendedForwards
        if self.healed:
            out += ", healed"
        return out


class Matcher:
    def __init__(self, textObjA, textObjB, threshold=THRESHOLD, cutoff=CUTOFF, ngramSize=NGRAM, removeStopwords=True, minDistance=DISTANCE, silent=True):
        print(f"Comparing: {textObjA.label} with {textObjB.label}")
        self.threshold = threshold
        self.ngramSize = ngramSize
        self.minDistance = minDistance
        self.silent = silent

        self.textA = textObjA
        self.textB = textObjB

        self.textAgrams = self.textA.ngrams(ngramSize)
        self.textBgrams = self.textB.ngrams(ngramSize)

        self.locationsA = []
        self.locationsB = []
        self.match_texts = []

        self.initial_matches = self.get_initial_matches()
        print(f"Found {len(self.initial_matches)} initial matches above threshold {threshold}")
        
        self.healed_matches = self.heal_neighboring_matches()
        print(f"After healing: {len(self.healed_matches)} match groups")
        
        self.extended_matches = self.extend_matches()
        self.extended_matches = [match for match in self.extended_matches
                                if min(match.sizeA, match.sizeB) >= cutoff]
        print(f"After extending and applying cutoff {cutoff}: {len(self.extended_matches)} significant matches")

        self.numMatches = len(self.extended_matches)
        self.similarity_score = self.calculate_similarity()
        print(f"Similarity score: {self.similarity_score}%")

    def calculate_similarity(self):
        if not self.extended_matches:
            return 0.0
            
        total_matched_tokens_A = sum(match.sizeA for match in self.extended_matches)
        total_tokens_A = len(self.textA.tokens)
        
        if total_tokens_A == 0:
            return 0.0
            
        similarity = (total_matched_tokens_A / total_tokens_A) * 100
        return round(similarity, 2)

    def get_initial_matches(self):
        sequence = SequenceMatcher(None, self.textAgrams, self.textBgrams)
        matchingBlocks = sequence.get_matching_blocks()
        highMatchingBlocks = [match for match in matchingBlocks if match.size > self.threshold]
        return highMatchingBlocks

    def getTokensText(self, text, start, length):
        if start < 0:
            start = 0
        
        matchTokens = text.tokens[start:start + length]
        
        if start >= len(text.spans) or start + length > len(text.spans):
            return ""
            
        spans = text.spans[start:start + length]
        if len(spans) == 0:
            passage = ""
        else:
            passage = text.text[spans[0][0]:spans[-1][-1]]
        return passage

    def getLocations(self, text, start, length, asPercentages=False):
        if start >= len(text.spans) or start + length > len(text.spans):
            return None
            
        spans = text.spans[start:start + length]
        if len(spans) == 0:
            return None
            
        if asPercentages:
            locations = (spans[0][0] / text.length, spans[-1][-1] / text.length)
        else:
            try:
                locations = (spans[0][0], spans[-1][-1])
            except IndexError:
                return None
        return locations

    def heal_neighboring_matches(self):
        healedMatches = []
        ignoreNext = False
        matches = self.initial_matches.copy()
        
        if len(matches) == 1:
            match = matches[0]
            sizeA, sizeB = match.size, match.size
            match = ExtendedMatch(match.a, match.b, sizeA, sizeB)
            healedMatches.append(match)
            return healedMatches
            
        for i, match in enumerate(matches):
            if i + 1 > len(matches) - 1:
                break
            nextMatch = matches[i + 1]
            
            if ignoreNext:
                ignoreNext = False
                continue
            else:
                if (nextMatch.a - (match.a + match.size)) < self.minDistance:
                    sizeA = (nextMatch.a + nextMatch.size) - match.a
                    sizeB = (nextMatch.b + nextMatch.size) - match.b
                    healed = ExtendedMatch(match.a, match.b, sizeA, sizeB)
                    healed.healed = True
                    healedMatches.append(healed)
                    ignoreNext = True
                else:
                    sizeA, sizeB = match.size, match.size
                    match = ExtendedMatch(match.a, match.b, sizeA, sizeB)
                    healedMatches.append(match)
        return healedMatches

    def edit_ratio(self, wordA, wordB):
        distance = editDistance(wordA, wordB)
        averageLength = (len(wordA) + len(wordB)) / 2
        return distance / averageLength

    def extend_matches(self, cutoff=0.4):
        extended = False
        for match in self.healed_matches:
            if match.a > 0 and match.b > 0 and len(self.textAgrams) > match.a - 1 and len(self.textBgrams) > match.b - 1:
                wordA = self.textAgrams[(match.a - 1)][0]
                wordB = self.textBgrams[(match.b - 1)][0]
                if self.edit_ratio(wordA, wordB) < cutoff:
                    match.a -= 1
                    match.b -= 1
                    match.sizeA += 1
                    match.sizeB += 1
                    match.extendedBackwards += 1
                    extended = True
                    
            idxA = match.a + match.sizeA + 1
            idxB = match.b + match.sizeB + 1
            if idxA >= len(self.textAgrams) or idxB >= len(self.textBgrams):
                continue
                
            wordA = self.textAgrams[idxA][-1] if idxA < len(self.textAgrams) else ""
            wordB = self.textBgrams[idxB][-1] if idxB < len(self.textBgrams) else ""
            
            if wordA and wordB and self.edit_ratio(wordA, wordB) < cutoff:
                match.sizeA += 1
                match.sizeB += 1
                match.extendedForwards += 1
                extended = True

        if extended:
            self.extend_matches()

        return self.healed_matches

    def match(self):
        matches_info = []
        for num, match in enumerate(self.extended_matches):
            lengthA = match.sizeA + self.ngramSize - 1
            matched_text = self.getTokensText(self.textA, match.a, lengthA)
            
            if matched_text:
                matches_info.append({
                    "source": self.textB.label,
                    "similarity": self.similarity_score,
                    "matched_text": matched_text[:200] + "..." if len(matched_text) > 200 else matched_text
                })

        return self.numMatches, matches_info, self.similarity_score


class AIDetector:
    """Client that calls an external service to get AI-generated score for text."""

    def __init__(self, endpoint: str = AI_SCORE_API_URL, timeout_sec: float = 15.0):
        self.endpoint = endpoint
        self.timeout = timeout_sec

    def analyze_text(self, text):
        """Send text to external API and return a normalized analysis dict.

        Expected external API response schema (from e5-small-lora API):
            { "ai_score": <float in [0,1]> }
        Returns a structure compatible with existing UI consumers:
            { ai_percentage, interpretation, chunks_analyzed, confidence }
        """
        print("Requesting AI-generated score from external service...")
        ai_score = 0.0
        try:
            r = requests.post(self.endpoint, json={"text": text}, timeout=self.timeout)
            r.raise_for_status()
            data = r.json()
            ai_score = float(data.get("ai_score", 0.0))
        except Exception as e:
            logging.warning(f"AI score API call failed: {e}")
            ai_score = 0.0

        ai_percentage = max(0.0, min(ai_score * 100.0, 100.0))

        # Simple interpretation mapping
        if ai_percentage < 15:
            interpretation = "Low AI probability"
        elif ai_percentage < 35:
            interpretation = "Moderate AI probability"
        elif ai_percentage < 60:
            interpretation = "High AI probability"
        else:
            interpretation = "Very high AI probability"

        # Single-shot classification; keep metadata stable
        result = {
            "ai_percentage": round(ai_percentage, 1),
            "interpretation": interpretation,
            "chunks_analyzed": 1,
            "confidence": 0.8,
        }
        print(f"AI analysis complete via external service: {result['ai_percentage']}%")
        return result


class PlagiarismAPI:
    def __init__(self):
        print("Initializing Plagiarism Detection API...")
        self.ai_detector = AIDetector()
        self.pdf_processor = PDFProcessor()
        self.ensure_nltk_data()
        
    def ensure_nltk_data(self):
        """Download required NLTK data"""
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            print("Downloading NLTK punkt tokenizers...")
            nltk.download('punkt')
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            print("Downloading NLTK stopwords...")
            nltk.download('stopwords')
    
    def download_and_process_file(self, url):
        """Download file from URL and extract text content, saving both original and text files"""
        try:
            print(f"Downloading file from: {url}")
            response = requests.get(url, timeout=60)
            response.raise_for_status()
            
            # Create unique download ID
            download_id = hashlib.md5(url.encode()).hexdigest()[:10]
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            download_folder = os.path.join(DOWNLOAD_DIR, f"{download_id}_{timestamp}")
            os.makedirs(download_folder, exist_ok=True)
            
            # Log download details
            download_log = {
                "url": url,
                "download_id": download_id,
                "timestamp": timestamp,
                "download_folder": download_folder,
                "response_status": response.status_code,
                "content_type": response.headers.get('content-type', 'unknown'),
                "content_length": len(response.content)
            }
            
            logging.info(f"Download details: {json.dumps(download_log)}")
            print(f"Downloaded {download_log['content_length']} bytes of type {download_log['content_type']}")
            
            # Determine filename from URL or Content-Disposition header
            filename = None
            if 'content-disposition' in response.headers:
                import re
                cd = response.headers['content-disposition']
                filename_match = re.findall('filename=(.+)', cd)
                if filename_match:
                    filename = filename_match[0].strip('"')
            
            if not filename:
                filename = os.path.basename(urlparse(url).path) or "document"
                
            # Ensure filename has proper extension based on content type
            content_type = response.headers.get('content-type', '').lower()
            if not os.path.splitext(filename)[1]:
                if 'pdf' in content_type:
                    filename += '.pdf'
                elif 'word' in content_type or 'docx' in content_type:
                    filename += '.docx'
                elif 'powerpoint' in content_type or 'pptx' in content_type:
                    filename += '.pptx'
                else:
                    filename += '.txt'
            
            # Secure filename
            safe_filename = secure_filename(filename)
            original_file_path = os.path.join(download_folder, safe_filename)
            
            # Save original file
            with open(original_file_path, 'wb') as f:
                f.write(response.content)
                
            print(f"Original file saved to: {original_file_path}")
                
            # Extract text using PDF processor
            print(f"Extracting text from downloaded file: {safe_filename}")
            extracted_text = self.pdf_processor.process_file_content(original_file_path, content_type)
            
            # Save extracted text to a .txt file
            text_filename = os.path.splitext(safe_filename)[0] + '.txt'
            text_file_path = os.path.join(TEXT_DIR, text_filename)
            
            with open(text_file_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
                
            print(f"Extracted text saved to: {text_file_path}")
                
            # Update log with extraction details
            extraction_log = {
                "original_file": original_file_path,
                "text_file": text_file_path,
                "text_length": len(extracted_text),
                "extraction_timestamp": datetime.now().isoformat()
            }
            
            logging.info(f"Extraction details: {json.dumps(extraction_log)}")
            
            return extracted_text, filename, text_file_path
            
        except Exception as e:
            logging.error(f"Error downloading and processing file from {url}: {e}\n{traceback.format_exc()}")
            print(f"Error downloading and processing file from {url}: {e}")
            raise
    
    def find_text_files_in_folder(self, target_folder):
        """Find all text files in the target folder for comparison"""
        if not os.path.exists(target_folder):
            return []
            
        text_files = []
        for file in os.listdir(target_folder):
            if file.lower().endswith('.txt'):
                file_path = os.path.join(target_folder, file)
                if os.path.isfile(file_path):
                    text_files.append(file_path)
                    
        return text_files
    
    def analyze_similarity_with_folder(self, main_text_path, main_text_content, filename):
        """Analyze text for similarity against all text files in the same folder"""
        print(f"Performing folder-based similarity analysis for: {filename}")
        
        # Create Text object for main text
        main_text = Text(main_text_content, filename, main_text_path)
        
        # Find other text files in the TEXT_DIR folder
        comparison_files = self.find_text_files_in_folder(TEXT_DIR)
        print(f"Found {len(comparison_files)} text files for comparison")
        
        # Remove the main text from comparison if it exists in the folder
        if main_text_path in comparison_files:
            comparison_files.remove(main_text_path)
            
        detailed_results = []
        max_similarity = 0.0
        max_similarity_source = ""
        total_comparisons = 0
        all_comparisons = []
        
        # Create a report folder
        report_timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        report_folder = os.path.join(REPORTS_DIR, f"report_{report_timestamp}")
        os.makedirs(report_folder, exist_ok=True)
        
        # Detailed comparison report
        comparison_report = []
        
        # Compare with each file in the folder
        for source_path in comparison_files:
            try:
                source_name = os.path.basename(source_path)
                print(f"Comparing with: {source_name}")
                
                # Read source content
                with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
                    source_content = f.read()
                
                # Create Text object for source
                source_text = Text(source_content, source_name, source_path)
                
                # Compare texts
                matcher = Matcher(main_text, source_text, 
                                 threshold=THRESHOLD, 
                                 cutoff=CUTOFF, 
                                 ngramSize=NGRAM, 
                                 minDistance=DISTANCE, 
                                 silent=True)
                
                num_matches, matches_info, similarity = matcher.match()
                
                # Add to comparison results
                comparison = {
                    "source": source_name,
                    "source_path": source_path,
                    "similarity": similarity,
                    "num_matches": num_matches
                }
                all_comparisons.append(comparison)
                
                # Add detailed info if there's significant similarity
                if similarity > 0:
                    for match in matches_info:
                        detailed_results.append(match)
                    
                    # Update max similarity if this is higher
                    if similarity > max_similarity:
                        max_similarity = similarity
                        max_similarity_source = source_name
                    
                    # Add to report
                    comparison_report.append({
                        "source": source_name,
                        "similarity": similarity,
                        "matches": matches_info
                    })
                
                total_comparisons += 1
                
            except Exception as e:
                logging.error(f"Error comparing with {source_path}: {e}\n{traceback.format_exc()}")
                print(f"Error comparing with {source_path}: {e}")
                continue
        
        # Save comparison report to JSON file
        report_file = os.path.join(report_folder, "similarity_report.json")
        with open(report_file, 'w', encoding='utf-8') as f:
            json.dump({
                "filename": filename,
                "timestamp": datetime.now().isoformat(),
                "max_similarity": max_similarity,
                "max_similarity_source": max_similarity_source,
                "total_comparisons": total_comparisons,
                "all_comparisons": all_comparisons,
                "detailed_comparisons": comparison_report
            }, f, indent=2)
            
        print(f"Similarity analysis complete. Max similarity: {max_similarity}% with {max_similarity_source}")
        print(f"Detailed report saved to: {report_file}")
        
        return {
            "similarity_score": round(max_similarity, 1),
            "max_similarity_source": max_similarity_source,
            "total_comparisons": total_comparisons,
            "detailed_results": detailed_results[:10],  # Limit to top 10 results
            "report_path": report_file
        }
    
    def calculate_plagiarism_score(self, similarity_score, ai_percentage):
        """Calculate overall plagiarism score"""
        print("Calculating overall plagiarism score...")
        # Weight the scores (you can adjust these weights)
        similarity_weight = 0.6
        ai_weight = 0.4
        
        plagiarism_score = (similarity_score * similarity_weight) + (ai_percentage * ai_weight)
        
        # Determine risk level
        if plagiarism_score < 15:
            risk_level = "low"
        elif plagiarism_score < 35:
            risk_level = "moderate"
        elif plagiarism_score < 60:
            risk_level = "high"
        else:
            risk_level = "critical"
            
        # Determine contributing factors
        contributing_factors = []
        if similarity_score > 25:
            contributing_factors.append("high_similarity_matches")
        if ai_percentage > 25:
            contributing_factors.append("ai_generated_content")
        if similarity_score > 15 and ai_percentage > 15:
            contributing_factors.append("multiple_detection_methods")
        if not contributing_factors:
            contributing_factors.append("low_risk_indicators")
        
        print(f"Plagiarism score: {round(plagiarism_score, 1)}%, Risk level: {risk_level}")    
        return {
            "plagiarism_score": round(plagiarism_score, 1),
            "risk_level": risk_level,
            "contributing_factors": contributing_factors
        }
    
    def generate_html_report(self, student_data, similarity_analysis, ai_analysis, plagiarism_analysis, content):
        """Generate comprehensive HTML report"""
        print("Generating HTML report...")
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # Determine color classes based on scores
        def get_score_class(score):
            if score < 15:
                return "low-risk"
            elif score < 35:
                return "medium-risk"
            elif score < 60:
                return "high-risk"
            else:
                return "critical-risk"
        
        html_report = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Plagiarism Analysis Report - {student_data['student_name']}</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 20px;
            background-color: #f8f9fa;
            color: #333;
            line-height: 1.6;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            overflow: hidden;
        }}
        .header {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 30px;
            text-align: center;
        }}
        .header h1 {{
            margin: 0;
            font-size: 28px;
            font-weight: 300;
        }}
        .content {{
            padding: 30px;
        }}
        .info-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }}
        .info-card {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            border-left: 4px solid #667eea;
        }}
        .info-card h3 {{
            margin: 0 0 10px 0;
            color: #667eea;
            font-size: 14px;
            font-weight: 600;
            text-transform: uppercase;
            letter-spacing: 1px;
        }}
        .info-card p {{
            margin: 0;
            font-size: 16px;
            font-weight: 500;
        }}
        .scores-section {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 30px;
            margin: 30px 0;
        }}
        .score-card {{
            background: white;
            border-radius: 10px;
            padding: 25px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
            border-top: 4px solid;
        }}
        .score-card.similarity {{
            border-top-color: #3498db;
        }}
        .score-card.ai-detection {{
            border-top-color: #e74c3c;
        }}
        .score-card.plagiarism {{
            border-top-color: #f39c12;
        }}
        .score-card h3 {{
            margin: 0 0 15px 0;
            font-size: 18px;
            font-weight: 600;
        }}
        .score-display {{
            font-size: 36px;
            font-weight: bold;
            margin: 10px 0;
        }}
        .low-risk {{ color: #27ae60; }}
        .medium-risk {{ color: #f39c12; }}
        .high-risk {{ color: #e74c3c; }}
        .critical-risk {{ color: #c0392b; }}
        .analysis-details {{
            margin-top: 30px;
        }}
        .details-card {{
            background: #f8f9fa;
            border-radius: 8px;
            padding: 20px;
            margin-bottom: 20px;
        }}
        .details-card h4 {{
            margin: 0 0 15px 0;
            color: #2c3e50;
        }}
        .matches-list {{
            list-style: none;
            padding: 0;
        }}
        .matches-list li {{
            background: white;
            padding: 15px;
            margin-bottom: 10px;
            border-radius: 5px;
            border-left: 3px solid #3498db;
        }}
        .risk-factors {{
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            margin-top: 10px;
        }}
        .risk-factor {{
            background: #e74c3c;
            color: white;
            padding: 5px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: 500;
        }}
        .footer {{
            background: #2c3e50;
            color: white;
            text-align: center;
            padding: 20px;
            font-size: 12px;
        }}
        .interpretation {{
            font-style: italic;
            color: #666;
            margin-top: 5px;
        }}
        .processing-info {{
            background: #e8f5e8;
            border: 1px solid #c3e6c3;
            border-radius: 5px;
            padding: 15px;
            margin-bottom: 20px;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Plagiarism Analysis Report</h1>
            <p>Comprehensive academic integrity assessment with PDF/OCR processing</p>
        </div>
        
        <div class="content">
            <div class="processing-info">
                <h4>Document Processing Information</h4>
                <p>This document was automatically processed using advanced PDF extraction and OCR technology. 
                Text content includes both direct text extraction and OCR results from images within the document.</p>
                <p><strong>Content Length:</strong> {len(content)} characters | 
                <strong>Word Count:</strong> {len(content.split())} words</p>
                <p><strong>Files compared:</strong> {similarity_analysis['total_comparisons']} text files from the document corpus</p>
            </div>
            
            <div class="info-grid">
                <div class="info-card">
                    <h3>Student</h3>
                    <p>{student_data['student_name']}</p>
                </div>
                <div class="info-card">
                    <h3>Classroom</h3>
                    <p>{student_data['classroom_name']}</p>
                </div>
                <div class="info-card">
                    <h3>Assignment ID</h3>
                    <p>{student_data['assignment_id']}</p>
                </div>
                <div class="info-card">
                    <h3>Analysis Date</h3>
                    <p>{timestamp}</p>
                </div>
            </div>
            
            <div class="scores-section">
                <div class="score-card similarity">
                    <h3>Similarity Analysis</h3>
                    <div class="score-display {get_score_class(similarity_analysis['similarity_score'])}">
                        {similarity_analysis['similarity_score']}%
                    </div>
                    <p>Compared against {similarity_analysis['total_comparisons']} documents</p>
                    {f'<p>Highest match with: {similarity_analysis["max_similarity_source"]}</p>' if similarity_analysis.get("max_similarity_source") else ''}
                </div>
                
                <div class="score-card ai-detection">
                    <h3>AI Detection</h3>
                    <div class="score-display {get_score_class(ai_analysis['ai_percentage'])}">
                        {ai_analysis['ai_percentage']}%
                    </div>
                    <p class="interpretation">{ai_analysis['interpretation']}</p>
                    <p>Confidence: {ai_analysis['confidence']}</p>
                </div>
                
                <div class="score-card plagiarism">
                    <h3>Overall Plagiarism Score</h3>
                    <div class="score-display {get_score_class(plagiarism_analysis['plagiarism_score'])}">
                        {plagiarism_analysis['plagiarism_score']}%
                    </div>
                    <p>Risk Level: <strong>{plagiarism_analysis['risk_level'].upper()}</strong></p>
                    <div class="risk-factors">
                        {' '.join([f'<span class="risk-factor">{factor.replace("_", " ").title()}</span>' for factor in plagiarism_analysis['contributing_factors']])}
                    </div>
                </div>
            </div>
            
            <div class="analysis-details">
                {f'''
                <div class="details-card">
                    <h4>Similarity Matches Found</h4>
                    <ul class="matches-list">
                        {"".join([f"<li><strong>{match['source']}</strong><br>Similarity: {match['similarity']}%<br><em>{match['matched_text']}</em></li>" for match in similarity_analysis['detailed_results'][:5]])}
                    </ul>
                </div>
                ''' if similarity_analysis['detailed_results'] else '<div class="details-card"><h4>No Significant Similarity Matches Found</h4><p>The document appears to have original content with no substantial matches in our database.</p></div>'}
                
                <div class="details-card">
                    <h4>AI Detection Analysis</h4>
                    <p>Analyzed {ai_analysis['chunks_analyzed']} text segments with {ai_analysis['confidence']} confidence level.</p>
                    <p><strong>Interpretation:</strong> {ai_analysis['interpretation']}</p>
                </div>
                
                <div class="details-card">
                    <h4>Processing Information</h4>
                    <p>All extracted text content has been saved for future comparison. A detailed report has been generated.</p>
                    <p>Similarity analysis parameters: Threshold={THRESHOLD}, Cutoff={CUTOFF}, N-gram Size={NGRAM}, Distance={DISTANCE}</p>
                </div>
            </div>
        </div>
        
        <div class="footer">
            <p>Generated by Enhanced Academic Integrity Analysis System with PDF/OCR Processing | {timestamp}</p>
            <p>This report is confidential and intended solely for academic evaluation purposes.</p>
        </div>
    </div>
</body>
</html>"""
        
        # Save the HTML report to a file
        report_path = os.path.join(REPORTS_DIR, f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write(html_report)
            
        print(f"HTML report generated and saved to: {report_path}")
        return html_report
    
    def log_user_score(self, student_id, scores):
        """Log user scores for tracking"""
        global user_scores_log
        
        if student_id not in user_scores_log:
            user_scores_log[student_id] = []
            
        user_scores_log[student_id].append({
            "timestamp": datetime.now().isoformat(),
            "scores": scores
        })
        
        # Keep only last 10 entries per user
        user_scores_log[student_id] = user_scores_log[student_id][-10:]
        
        print(f"User scores logged for student ID: {student_id}")
    
    def analyze_document(self, payload):
        """Main analysis function with integrated PDF/OCR processing and folder-based similarity"""
        try:
            print(f"Starting analysis for document submitted by: {payload['student_name']}")
            
            # Download and process file with enhanced PDF/OCR capabilities
            content, filename, text_file_path = self.download_and_process_file(payload['file_url'])
            
            # Validate that we got content
            if not content or len(content.strip()) < 10:
                error_msg = "Could not extract meaningful content from the document"
                logging.error(error_msg)
                print(error_msg)
                raise ValueError(error_msg)
            
            # Perform folder-based similarity analysis
            similarity_analysis = self.analyze_similarity_with_folder(text_file_path, content, payload['student_name'])
            
            # Perform AI detection analysis
            ai_analysis = self.ai_detector.analyze_text(content)
            
            # Calculate overall plagiarism score
            plagiarism_analysis = self.calculate_plagiarism_score(
                similarity_analysis['similarity_score'],
                ai_analysis['ai_percentage']
            )
            
            # Generate HTML report
            html_report = self.generate_html_report(
                payload, similarity_analysis, ai_analysis, plagiarism_analysis, content
            )
            
            # Log scores
            scores = {
                "similarity_score": similarity_analysis['similarity_score'],
                "ai_percentage": ai_analysis['ai_percentage'],
                "plagiarism_score": plagiarism_analysis['plagiarism_score'],
                "content_length": len(content),
                "word_count": len(content.split()),
                "text_file": text_file_path
            }
            self.log_user_score(payload['student_id'], scores)
            
            print(f"Analysis completed successfully for: {payload['student_name']}")
            
            return {
                "status": "completed",
                "similarity_analysis": similarity_analysis,
                "ai_analysis": ai_analysis,
                "plagiarism_analysis": plagiarism_analysis,
                "report_html": html_report,
                "processing_info": {
                    "filename": filename,
                    "content_length": len(content),
                    "word_count": len(content.split()),
                    "text_file": text_file_path
                },
                "errors": []
            }
            
        except Exception as e:
            logging.error(f"Error in analysis: {e}\n{traceback.format_exc()}")
            print(f"Error in analysis: {e}")
            return {
                "status": "failed",
                "similarity_analysis": {"similarity_score": 0, "total_comparisons": 0, "detailed_results": []},
                "ai_analysis": {"ai_percentage": 0, "interpretation": "Analysis failed", "chunks_analyzed": 0, "confidence": 0},
                "plagiarism_analysis": {"plagiarism_score": 0, "risk_level": "unknown", "contributing_factors": ["analysis_error"]},
                "report_html": f"<html><body><h1>Analysis Failed</h1><p>Error: {str(e)}</p></body></html>",
                "processing_info": {"filename": "unknown", "content_length": 0, "word_count": 0},
                "errors": [str(e)]
            }


# Initialize API instance
api = PlagiarismAPI()

@app.route('/grade/analyze', methods=['POST'])
def analyze_submission():
    """Main API endpoint for plagiarism analysis with PDF/OCR processing"""
    try:
        print("Received analyze submission request")
        # Validate request
        if not request.is_json:
            print("Error: Content-Type must be application/json")
            return jsonify({"error": "Content-Type must be application/json"}), 400
            
        payload = request.get_json()
        
        # Validate required fields
        required_fields = ['student_id', 'student_name', 'file_url', 'assignment_id', 'classroom_name']
        for field in required_fields:
            if field not in payload:
                print(f"Error: Missing required field: {field}")
                return jsonify({"error": f"Missing required field: {field}"}), 400
        
        # Perform analysis with integrated PDF/OCR processing
        result = api.analyze_document(payload)
        
        return jsonify(result)
        
    except Exception as e:
        logging.error(f"API error: {e}\n{traceback.format_exc()}")
        print(f"API error: {e}")
        return jsonify({
            "status": "failed",
            "similarity_analysis": {"similarity_score": 0, "total_comparisons": 0, "detailed_results": []},
            "ai_analysis": {"ai_percentage": 0, "interpretation": "Analysis failed", "chunks_analyzed": 0, "confidence": 0},
            "plagiarism_analysis": {"plagiarism_score": 0, "risk_level": "unknown", "contributing_factors": ["api_error"]},
            "report_html": f"<html><body><h1>Analysis Failed</h1><p>API Error: {str(e)}</p></body></html>",
            "processing_info": {"filename": "unknown", "content_length": 0, "word_count": 0},
            "errors": [str(e)]
        }), 500


@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    print("Health check requested")
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "service": "Enhanced Plagiarism Detection API with PDF/OCR Processing",
        "capabilities": [
            "PDF text extraction",
            "OCR for image-based documents",
            "DOCX/PPTX processing",
            "Image preprocessing",
            "Multi-format support",
            "Folder-based similarity comparison",
            "Extracted text preservation"
        ],
        "config": {
            "threshold": THRESHOLD,
            "cutoff": CUTOFF,
            "ngram": NGRAM,
            "distance": DISTANCE
        }
    })


@app.route('/test-processing', methods=['POST'])
def test_file_processing():
    """Test endpoint for file processing capabilities"""
    try:
        print("Received test file processing request")
        if not request.is_json:
            print("Error: Content-Type must be application/json")
            return jsonify({"error": "Content-Type must be application/json"}), 400
            
        payload = request.get_json()
        
        if 'file_url' not in payload:
            print("Error: Missing required field: file_url")
            return jsonify({"error": "Missing required field: file_url"}), 400
        
        # Test file processing
        content, filename, text_file_path = api.download_and_process_file(payload['file_url'])
        
        print(f"Test processing successful for: {filename}")
        
        return jsonify({
            "status": "success",
            "filename": filename,
            "content_preview": content[:500] + "..." if len(content) > 500 else content,
            "content_length": len(content),
            "word_count": len(content.split()),
            "text_file_path": text_file_path
        })
        
    except Exception as e:
        logging.error(f"Test processing error: {e}\n{traceback.format_exc()}")
        print(f"Test processing error: {e}")
        return jsonify({
            "status": "failed",
            "error": str(e)
        }), 500


@app.route('/user/<student_id>/scores', methods=['GET'])
def get_user_scores(student_id):
    """Get historical scores for a user"""
    print(f"Retrieving score history for student ID: {student_id}")
    global user_scores_log
    
    if student_id in user_scores_log:
        return jsonify({
            "student_id": student_id,
            "scores_history": user_scores_log[student_id]
        })
    else:
        print(f"No score history found for student ID: {student_id}")
        return jsonify({
            "student_id": student_id,
            "scores_history": [],
            "message": "No scores found for this user"
        })


@app.route('/stats', methods=['GET'])
def get_stats():
    """Get system statistics"""
    print("Retrieving system statistics")
    global user_scores_log
    
    # Count total files in each directory
    download_count = len([f for f in os.listdir(DOWNLOAD_DIR) if os.path.isdir(os.path.join(DOWNLOAD_DIR, f))])
    text_count = len([f for f in os.listdir(TEXT_DIR) if f.endswith('.txt')])
    report_count = len([f for f in os.listdir(REPORTS_DIR) if f.endswith('.html') or f.endswith('.json')])
    
    total_users = len(user_scores_log)
    total_analyses = sum(len(scores) for scores in user_scores_log.values())
    
    print(f"Stats: {total_users} users, {total_analyses} analyses, {text_count} text files")
    
    return jsonify({
        "total_users_analyzed": total_users,
        "total_analyses_performed": total_analyses,
        "stored_files": {
            "downloads": download_count,
            "extracted_texts": text_count,
            "reports": report_count
        },
        "system_status": "operational",
        "configuration": {
            "threshold": THRESHOLD,
            "cutoff": CUTOFF,
            "ngram_size": NGRAM,
            "min_distance": DISTANCE
        },
        "features": [
            "PDF text extraction with PyMuPDF",
            "OCR processing with Tesseract",
            "Image preprocessing with OpenCV",
            "DOCX/PPTX support",
            "Multi-format document processing",
            "Folder-based similarity comparison",
            "Extracted text preservation for future analysis"
        ]
    })


@app.route('/compare-texts', methods=['POST'])
def compare_specific_texts():
    """Compare two specific text files and generate a similarity report"""
    try:
        print("Received text comparison request")
        if not request.is_json:
            print("Error: Content-Type must be application/json")
            return jsonify({"error": "Content-Type must be application/json"}), 400
            
        payload = request.get_json()
        
        # Validate required fields
        if 'text1_path' not in payload or 'text2_path' not in payload:
            print("Error: Missing required fields: text1_path and/or text2_path")
            return jsonify({"error": "Missing required fields: text1_path and text2_path"}), 400
        
        text1_path = payload['text1_path']
        text2_path = payload['text2_path']
        
        # Validate that files exist
        if not os.path.exists(text1_path) or not os.path.exists(text2_path):
            print("Error: One or both text files do not exist")
            return jsonify({"error": "One or both text files do not exist"}), 400
        
        # Read file contents
        try:
            with open(text1_path, 'r', encoding='utf-8', errors='replace') as f:
                text1_content = f.read()
                
            with open(text2_path, 'r', encoding='utf-8', errors='replace') as f:
                text2_content = f.read()
        except Exception as e:
            print(f"Error reading text files: {e}")
            return jsonify({"error": f"Error reading text files: {str(e)}"}), 500
        
        # Create Text objects
        text1_name = os.path.basename(text1_path)
        text2_name = os.path.basename(text2_path)
        
        text1_obj = Text(text1_content, text1_name, text1_path)
        text2_obj = Text(text2_content, text2_name, text2_path)
        
        # Compare texts
        matcher = Matcher(text1_obj, text2_obj, 
                        threshold=THRESHOLD, 
                        cutoff=CUTOFF, 
                        ngramSize=NGRAM, 
                        minDistance=DISTANCE, 
                        silent=True)
        
        num_matches, matches_info, similarity = matcher.match()
        
        # Create comparison report
        report = {
            "file1": text1_name,
            "file2": text2_name,
            "similarity_score": similarity,
            "num_matches": num_matches,
            "matches": matches_info if len(matches_info) <= 10 else matches_info[:10],
            "comparison_timestamp": datetime.now().isoformat(),
            "parameters": {
                "threshold": THRESHOLD,
                "cutoff": CUTOFF,
                "ngram": NGRAM,
                "distance": DISTANCE
            }
        }
        
        print(f"Text comparison complete: {similarity}% similarity between {text1_name} and {text2_name}")
        
        # Save report to file
        report_file = os.path.join(REPORTS_DIR, f"comparison_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json")
        with open(report_file, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)
        
        return jsonify({
            "status": "success",
            "similarity": similarity,
            "num_matches": num_matches,
            "file1": text1_name,
            "file2": text2_name,
            "report_file": report_file,
            "matches_preview": matches_info[:3] if matches_info else []
        })
        
    except Exception as e:
        logging.error(f"Text comparison error: {e}\n{traceback.format_exc()}")
        print(f"Text comparison error: {e}")
        return jsonify({
            "status": "failed",
            "error": str(e)
        }), 500


@app.route('/cleanup-old-files', methods=['POST'])
def cleanup_old_files():
    """Remove old files to free up disk space"""
    try:
        print("Received cleanup request")
        if not request.is_json:
            print("Error: Content-Type must be application/json")
            return jsonify({"error": "Content-Type must be application/json"}), 400
            
        payload = request.get_json()
        
        # Get days threshold from payload or use default (30 days)
        days_threshold = payload.get('days_threshold', 30)
        
        # Calculate cutoff timestamp
        cutoff_time = datetime.now().timestamp() - (days_threshold * 86400)  # 86400 seconds in a day
        
        removed_files = {
            "downloads": 0,
            "text_files": 0,
            "reports": 0
        }
        
        # Clean up download directory
        print(f"Cleaning up files older than {days_threshold} days...")
        for folder in os.listdir(DOWNLOAD_DIR):
            folder_path = os.path.join(DOWNLOAD_DIR, folder)
            if os.path.isdir(folder_path):
                if os.path.getctime(folder_path) < cutoff_time:
                    try:
                        shutil.rmtree(folder_path)
                        removed_files["downloads"] += 1
                    except Exception as e:
                        logging.error(f"Error removing folder {folder_path}: {e}")
        
        # Clean up text files
        for file in os.listdir(TEXT_DIR):
            file_path = os.path.join(TEXT_DIR, file)
            if os.path.isfile(file_path) and os.path.getctime(file_path) < cutoff_time:
                try:
                    os.remove(file_path)
                    removed_files["text_files"] += 1
                except Exception as e:
                    logging.error(f"Error removing text file {file_path}: {e}")
        
        # Clean up reports
        for file in os.listdir(REPORTS_DIR):
            file_path = os.path.join(REPORTS_DIR, file)
            if os.path.isfile(file_path) and os.path.getctime(file_path) < cutoff_time:
                try:
                    os.remove(file_path)
                    removed_files["reports"] += 1
                except Exception as e:
                    logging.error(f"Error removing report file {file_path}: {e}")
        
        print(f"Cleanup complete. Removed: {removed_files['downloads']} download folders, "
              f"{removed_files['text_files']} text files, {removed_files['reports']} reports")
              
        return jsonify({
            "status": "success",
            "days_threshold": days_threshold,
            "removed_files": removed_files,
            "timestamp": datetime.now().isoformat()
        })
        
    except Exception as e:
        logging.error(f"Cleanup error: {e}\n{traceback.format_exc()}")
        print(f"Cleanup error: {e}")
        return jsonify({
            "status": "failed",
            "error": str(e)
        }), 500


if __name__ == '__main__':
    # Ensure all required directories exist
    for directory in [DOWNLOAD_DIR, TEXT_DIR, REPORTS_DIR, LOGS_DIR]:
        os.makedirs(directory, exist_ok=True)
    
    print("\n" + "="*80)
    print("Starting Enhanced Plagiarism Detection API Service with PDF/OCR Processing...")
    print("="*80)
    print(f"Configuration: THRESHOLD={THRESHOLD}, CUTOFF={CUTOFF}, NGRAM={NGRAM}, DISTANCE={DISTANCE}")
    print(f"Server will run on http://localhost:5000")
    print("\nDirectories:")
    print(f"✓ Downloads: {os.path.abspath(DOWNLOAD_DIR)}")
    print(f"✓ Extracted Text: {os.path.abspath(TEXT_DIR)}")
    print(f"✓ Reports: {os.path.abspath(REPORTS_DIR)}")
    print(f"✓ Logs: {os.path.abspath(LOGS_DIR)}")
    print("\nFeatures enabled:")
    print("✓ PDF text extraction with PyMuPDF")
    print("✓ OCR processing with Tesseract")
    print("✓ Image preprocessing with OpenCV/PIL")
    print("✓ DOCX/PPTX document support")
    print("✓ Multi-format file processing")
    print("✓ Folder-based similarity comparison")
    print("✓ Text extraction preservation for future analysis")
    print("\nAvailable endpoints:")
    print("POST /grade/analyze - Main plagiarism analysis endpoint")
    print("POST /test-processing - Test file processing capabilities")
    print("GET /health - Health check with capabilities")
    print("GET /user/<student_id>/scores - Get user's score history")
    print("GET /stats - Get system statistics")
    print("POST /compare-texts - Compare two specific text files")
    print("POST /cleanup-old-files - Remove old files to free up disk space")
    
    # Test Tesseract installation on startup
    try:
        pytesseract.get_tesseract_version()
        print(f"✓ Tesseract OCR ready: {pytesseract.get_tesseract_version()}")
    except Exception as e:
        print(f"⚠ Tesseract warning: {e}")
        print("  OCR functionality may be limited")
    
    # Count existing files
    text_count = len([f for f in os.listdir(TEXT_DIR) if f.endswith('.txt')])
    if text_count > 0:
        print(f"✓ Found {text_count} existing text files for comparison")
    
    print("="*80)
    
    app.run(host='0.0.0.0', port=5000, debug=True)