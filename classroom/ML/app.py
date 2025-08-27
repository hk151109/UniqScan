from flask import Flask, request, jsonify
import os
import json
import logging
from logging.handlers import RotatingFileHandler
import nltk
import time
import random
import hashlib
import requests
import tempfile
from datetime import datetime
from difflib import SequenceMatcher
from nltk.metrics.distance import edit_distance as editDistance
from nltk.stem.lancaster import LancasterStemmer
from nltk.util import ngrams
from pathlib import Path
import re
from urllib.parse import urlparse
from flask import Flask, request, jsonify
import os
import logging
from logging.handlers import RotatingFileHandler
import nltk
import time
import random
import hashlib
import requests
import tempfile
from datetime import datetime
from difflib import SequenceMatcher
from nltk.metrics.distance import edit_distance as editDistance
from nltk.stem.lancaster import LancasterStemmer
from nltk.util import ngrams
import re
from urllib.parse import urlparse
from werkzeug.utils import secure_filename

# PDF/OCR Processing imports
import pytesseract
from PIL import Image, ImageEnhance
import fitz  # PyMuPDF for handling PDFs
import cv2
import numpy as np
import traceback
from docx import Document
from pptx import Presentation
import csv
import platform

# ---------------------------------------------
# App and Config
# ---------------------------------------------
app = Flask(__name__)

# Easy-tweak matching params
THRESHOLD = 3       # --threshold
CUTOFF = 5          # --cutoff
NGRAM = 3           # --ngram
MIN_DISTANCE = 8    # --distance

BASE_DIR = os.path.dirname(__file__)
CORPUS_DIR = os.path.join(BASE_DIR, 'corpus')
LOGS_DIR = os.path.join(BASE_DIR, 'logs')
LOG_FILE = os.path.join(LOGS_DIR, 'api.log')

os.makedirs(CORPUS_DIR, exist_ok=True)
from urllib.parse import urlparse

logger = logging.getLogger('plagiarism_api')
logger.setLevel(logging.DEBUG)
if not logger.handlers:
    fh = RotatingFileHandler(LOG_FILE, maxBytes=2_000_000, backupCount=5, encoding='utf-8')
    fh.setLevel(logging.DEBUG)
    fh.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(name)s - %(message)s'))
    logger.addHandler(fh)
    ch = logging.StreamHandler()
    ch.setLevel(logging.INFO)
    ch.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
    logger.addHandler(ch)

# Tesseract on Windows
if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"

# In-memory scores history
user_scores_log = {}

# ---------------------------------------------
# PDF/Image/Text processors
# ---------------------------------------------
class PDFProcessor:
    def __init__(self):
        self.check_tesseract_installation()

    def check_tesseract_installation(self):
        try:
            pytesseract.get_tesseract_version()
            logger.info(f"Tesseract version: {pytesseract.get_tesseract_version()}")
            return True
        except Exception as e:
            logger.warning(f"Tesseract not configured: {e}")
            return False

    def preprocess_image(self, image, q=1.5):
        try:
            enhanced = image.copy()
            enhanced = ImageEnhance.Contrast(enhanced).enhance(q)
            enhanced = ImageEnhance.Sharpness(enhanced).enhance(q)
            arr = np.array(enhanced)
            if len(arr.shape) == 2:
                arr = cv2.cvtColor(arr, cv2.COLOR_GRAY2RGB)
            gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
            blurred = cv2.GaussianBlur(gray, (3, 3), 0)
            kernel = np.ones((1, 1), np.uint8)
            morph = cv2.morphologyEx(blurred, cv2.MORPH_CLOSE, kernel)
            binary = cv2.adaptiveThreshold(morph, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_, cv2.THRESH_BINARY, 11, 2)
            binary = cv2.dilate(binary, kernel, iterations=1)
            return enhanced, Image.fromarray(binary)
        except Exception as e:
            logger.error(f"preprocess_image error: {e}")
            return image, image

    def perform_ocr(self, image, lang="eng", config='--psm 6'):
        try:
            _, bin_img = self.preprocess_image(image)
            ocr = pytesseract.image_to_data(bin_img, lang=lang, config=config, output_type=pytesseract.Output.DICT)
            parts = []
            for i, conf in enumerate(ocr['conf']):
                try:
                    if float(conf) > 30 and ocr['text'][i].strip():
                        parts.append(ocr['text'][i])
                except Exception:
                    continue
            text = re.sub(r'\s+', ' ', ' '.join(parts)).strip()
            return text
        except Exception as e:
            logger.error(f"OCR error: {e}")
            return ""

    def extract_text_from_pdf(self, file_path):
        try:
            doc = fitz.open(file_path)
            out = []
            for idx in range(len(doc)):
                page = doc.load_page(idx)
                text = page.get_text()
                if len(text.strip()) > 50:
                    out.append(f"\n\n--- Page {idx+1} ---\n\n{text}")
                else:
                    try:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        ocr_text = self.perform_ocr(img)
                        if ocr_text.strip():
                            out.append(f"\n\n--- Page {idx+1} (OCR) ---\n\n{ocr_text}")
                    except Exception as e:
                        logger.error(f"Page OCR error p{idx}: {e}")
                for j, info in enumerate(page.get_images(full=True)):
                    try:
                        xref = info[0]
                        base = fitz.Pixmap(doc, xref)
                        if base.n < 5:
                            pil = Image.frombytes("RGB", [base.width, base.height], base.samples)
                        else:
                            cmyk = fitz.Pixmap(fitz.csRGB, base)
                            pil = Image.frombytes("RGB", [cmyk.width, cmyk.height], cmyk.samples)
                        if pil.width >= 50 and pil.height >= 50:
                            text_i = self.perform_ocr(pil)
                            if text_i.strip():
                                out.append(f"\n\n--- Page {idx+1}, Image {j+1} (OCR) ---\n\n{text_i}")
                    except Exception as e:
                        logger.error(f"Image OCR error p{idx}i{j}: {e}")
                        continue
            doc.close()
            return '\n'.join(out).strip()
        except Exception as e:
            logger.error(f"PDF extract error: {e}\n{traceback.format_exc()}")
            return f"Error extracting content from PDF: {str(e)}"

    def extract_text_from_docx(self, file_path):
        try:
            doc = Document(file_path)
            segs = []
            for p in doc.paragraphs:
                if p.text.strip():
                    segs.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        segs.append(' | '.join(cells))
            return '\n\n'.join(segs)
        except Exception as e:
            logger.error(f"DOCX extract error: {e}")
            return f"Error extracting content from DOCX: {str(e)}"

    def extract_text_from_pptx(self, file_path):
        try:
            prs = Presentation(file_path)
            segs = []
            for s_idx, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {s_idx+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if len(slide_text) > 1:
                    segs.append('\n'.join(slide_text))
            return '\n\n'.join(segs)
        except Exception as e:
            logger.error(f"PPTX extract error: {e}")
            return f"Error extracting content from PPTX: {str(e)}"

    def extract_text_from_txt(self, file_path):
        try:
            for enc in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=enc) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            with open(file_path, 'rb') as f:
                return f.read().decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"TXT read error: {e}")
            return f"Error reading text file: {str(e)}"

    def extract_text_from_csv(self, file_path):
        try:
            segs = []
            for enc in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=enc) as f:
                        for i, row in enumerate(csv.reader(f)):
                            if i > 100:
                                segs.append('...truncated, showing first 100 rows...')
                                break
                            segs.append(' '.join(row))
                    return '\n'.join(segs)
                except UnicodeDecodeError:
                    continue
            return 'Error: Could not decode CSV file'
        except Exception as e:
            logger.error(f"CSV read error: {e}")
            return f"Error reading CSV file: {str(e)}"

    def extract_text_from_image(self, file_path):
        try:
            img = Image.open(file_path)
            text = self.perform_ocr(img)
            return text if text.strip() else 'No text detected in image'
        except Exception as e:
            logger.error(f"Image process error: {e}")
            return f"Error processing image file: {str(e)}"

    def process_file_content(self, file_path, content_type=None):
        try:
            if not content_type:
                _, ext = os.path.splitext(file_path)
                ext = ext.lower().lstrip('.')
            else:
                if '/' in content_type:
                    ext = content_type.split('/')[-1]
                else:
                    _, ext = os.path.splitext(file_path)
                    ext = ext.lower().lstrip('.')
            if ext == 'pdf':
                return self.extract_text_from_pdf(file_path)
            if ext in ['docx', 'doc']:
                return self.extract_text_from_docx(file_path)
            if ext in ['pptx', 'ppt']:
                return self.extract_text_from_pptx(file_path)
            if ext == 'txt':
                return self.extract_text_from_txt(file_path)
            if ext == 'csv':
                return self.extract_text_from_csv(file_path)
            if ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif']:
                return self.extract_text_from_image(file_path)
            logger.warning(f"Unknown file type: {ext}. Attempting text extraction as txt.")
            return self.extract_text_from_txt(file_path)
        except Exception as e:
            logger.error(f"Process file error: {e}")
            return f"Error processing file: {str(e)}"

# ---------------------------------------------
# Tokenization and matching
# ---------------------------------------------
class Text:
    def __init__(self, raw_text, label, filepath, removeStopwords=True):
        self.text = ' \n '.join(raw_text) if isinstance(raw_text, list) else raw_text
        self.label = label
        self.filepath = filepath
        self.preprocess(self.text)
        self.tokens = self.getTokens(removeStopwords)
        self.trigrams = self.ngrams(3)
        self.checksum = hashlib.md5(self.text.encode()).hexdigest()

    def preprocess(self, text):
        self.text = re.sub(r'([A-Za-z])- ([a-z])', r'\1\2', text)

    def getTokens(self, removeStopwords=True):
        tokenizer = nltk.RegexpTokenizer('[a-zA-Z]\\w+\'?\\w*')
        spans = list(tokenizer.span_tokenize(self.text))
        self.length = spans[-1][-1] if spans else 0
        tokens = [t.lower() for t in tokenizer.tokenize(self.text)]
        stemmer = LancasterStemmer()
        tokens = [stemmer.stem(t) for t in tokens]
        if not removeStopwords:
            self.spans = spans
            return tokens
        stopwords = set(nltk.corpus.stopwords.words('english'))
        tokenSpans = [(t, s) for t, s in zip(tokens, spans) if t not in stopwords]
        self.spans = [s for _, s in tokenSpans]
        return [t for t, _ in tokenSpans]

    def ngrams(self, n):
        return list(ngrams(self.tokens, n))

class ExtendedMatch:
    def __init__(self, a, b, sizeA, sizeB):
        self.a = a
        self.b = b
        self.sizeA = sizeA
        self.sizeB = sizeB
        self.healed = False
        self.extendedBackwards = 0
        self.extendedForwards = 0

class Matcher:
    def __init__(self, textObjA, textObjB, threshold=3, cutoff=5, ngramSize=3, removeStopwords=True, minDistance=8, silent=True):
        self.threshold = threshold
        self.ngramSize = ngramSize
        self.minDistance = minDistance
        self.silent = silent
        self.textA = textObjA
        self.textB = textObjB
        self.textAgrams = self.textA.ngrams(ngramSize)
        self.textBgrams = self.textB.ngrams(ngramSize)
        self.initial_matches = self.get_initial_matches()
        self.healed_matches = self.heal_neighboring_matches()
        self.extended_matches = self.extend_matches()
        self.extended_matches = [m for m in self.extended_matches if min(m.sizeA, m.sizeB) >= cutoff]
        self.numMatches = len(self.extended_matches)
        self.similarity_score = self.calculate_similarity()

    def calculate_similarity(self):
        if not self.extended_matches:
            return 0.0
        total_matched_tokens_A = sum(m.sizeA for m in self.extended_matches)
        total_tokens_A = len(self.textA.tokens)
        if total_tokens_A == 0:
            return 0.0
        return round((total_matched_tokens_A / total_tokens_A) * 100, 2)

    def get_initial_matches(self):
        seq = SequenceMatcher(None, self.textAgrams, self.textBgrams)
        return [m for m in seq.get_matching_blocks() if m.size > self.threshold]

    def getTokensText(self, text, start, length):
        start = max(0, start)
        if start >= len(text.spans) or start + length > len(text.spans):
            return ""
        spans = text.spans[start:start + length]
        return text.text[spans[0][0]:spans[-1][-1]] if spans else ""

    def heal_neighboring_matches(self):
        healed = []
        matches = self.initial_matches.copy()
        if len(matches) == 1:
            m = matches[0]
            healed.append(ExtendedMatch(m.a, m.b, m.size, m.size))
            return healed
        ignoreNext = False
        for i, m in enumerate(matches):
            if i + 1 > len(matches) - 1:
                break
            n = matches[i + 1]
            if ignoreNext:
                ignoreNext = False
                continue
            if (n.a - (m.a + m.size)) < MIN_DISTANCE:
                sizeA = (n.a + n.size) - m.a
                sizeB = (n.b + n.size) - m.b
                h = ExtendedMatch(m.a, m.b, sizeA, sizeB)
                h.healed = True
                healed.append(h)
                ignoreNext = True
            else:
                healed.append(ExtendedMatch(m.a, m.b, m.size, m.size))
        return healed

    def edit_ratio(self, wordA, wordB):
        distance = editDistance(wordA, wordB)
        avg = (len(wordA) + len(wordB)) / 2
        return distance / avg if avg else 1.0

    def extend_matches(self, cutoff=0.4):
        extended = False
        for m in self.healed_matches:
            if m.a > 0 and m.b > 0 and len(self.textAgrams) > m.a - 1 and len(self.textBgrams) > m.b - 1:
                wordA = self.textAgrams[(m.a - 1)][0]
                wordB = self.textBgrams[(m.b - 1)][0]
                if self.edit_ratio(wordA, wordB) < cutoff:
                    m.a -= 1
                    m.b -= 1
                    m.sizeA += 1
                    m.sizeB += 1
                    m.extendedBackwards += 1
                    extended = True
            idxA = m.a + m.sizeA + 1
            idxB = m.b + m.sizeB + 1
            if idxA >= len(self.textAgrams) or idxB >= len(self.textBgrams):
                continue
            wordA = self.textAgrams[idxA][-1] if idxA < len(self.textAgrams) else ""
            wordB = self.textBgrams[idxB][-1] if idxB < len(self.textBgrams) else ""
            if wordA and wordB and self.edit_ratio(wordA, wordB) < cutoff:
                m.sizeA += 1
                m.sizeB += 1
                m.extendedForwards += 1
                extended = True
        if extended:
            self.extend_matches()
        return self.healed_matches

    def match(self):
        infos = []
        for m in self.extended_matches:
            lengthA = m.sizeA + self.ngramSize - 1
            txt = self.getTokensText(self.textA, m.a, lengthA)
            if txt:
                infos.append({
                    'source': self.textB.label,
                    'similarity': self.similarity_score,
                    'matched_text': txt[:200] + '...' if len(txt) > 200 else txt
                })
        return self.numMatches, infos, self.similarity_score

# ---------------------------------------------
# AI Detector (mock)
# ---------------------------------------------
class AIDetector:
    def analyze_text(self, text):
        wc = len(text.split())
        base = random.uniform(5.0, 35.0)
        if wc > 1000:
            base *= 0.8
        elif wc < 200:
            base *= 1.2
        ai = min(95.0, max(0.0, base))
        interp = (
            'Low AI probability' if ai < 15 else
            'Moderate AI probability' if ai < 35 else
            'High AI probability' if ai < 60 else
            'Very high AI probability'
        )
        conf = random.uniform(0.8, 0.95) if ai < 20 or ai > 70 else random.uniform(0.6, 0.85)
        return {
            'ai_percentage': round(ai, 1),
            'interpretation': interp,
            'chunks_analyzed': max(1, wc // 50),
            'confidence': round(conf, 2)
        }

# ---------------------------------------------
# Main API
# ---------------------------------------------
class PlagiarismAPI:
    def __init__(self):
        self.ai_detector = AIDetector()
        self.pdf_processor = PDFProcessor()
        self.ensure_nltk()

    def ensure_nltk(self):
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords')

    def download_and_process_file(self, url):
        try:
            resp = requests.get(url, timeout=60)
            resp.raise_for_status()
            temp_dir = tempfile.mkdtemp()
            filename = None
            if 'content-disposition' in resp.headers:
                cd = resp.headers['content-disposition']
                mm = re.findall('filename=(.+)', cd)
                if mm:
                    filename = mm[0].strip('"')
            if not filename:
                filename = os.path.basename(urlparse(url).path) or 'document'
            ctype = resp.headers.get('content-type', '').lower()
            if not os.path.splitext(filename)[1]:
                if 'pdf' in ctype:
                    filename += '.pdf'
                elif 'word' in ctype or 'docx' in ctype:
                    filename += '.docx'
                elif 'powerpoint' in ctype or 'pptx' in ctype:
                    filename += '.pptx'
                else:
                    filename += '.txt'
            temp_path = os.path.join(temp_dir, secure_filename(filename))
            logger.info(f"Downloading: {url}")
            with open(temp_path, 'wb') as f:
                f.write(resp.content)
            t0 = time.time()
            text = self.pdf_processor.process_file_content(temp_path, ctype)
            logger.info(f"Extraction OK in {time.time()-t0:.2f}s, len={len(text)}")
            try:
                os.unlink(temp_path)
                os.rmdir(temp_dir)
            except Exception:
                pass
            return text, filename
        except Exception as e:
            logger.error(f"download_and_process_file error: {e}")
            raise

    def save_text_to_corpus(self, text, original_filename, student_id, assignment_id):
        base = os.path.splitext(original_filename or 'document')[0]
        base = re.sub(r'[^A-Za-z0-9._-]+', '_', base)[:80]
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        checksum = hashlib.md5(text.encode('utf-8', errors='ignore')).hexdigest()[:8]
        sid = re.sub(r'[^A-Za-z0-9_-]+', '_', str(student_id))
        aid = re.sub(r'[^A-Za-z0-9_-]+', '_', str(assignment_id))
        fname = f"{base}__{sid}__{aid}__{ts}__{checksum}.txt"
        path = os.path.join(CORPUS_DIR, fname)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(text)
        logger.info(f"Saved text to corpus: {path}")
        return path

    def _read_text_file(self, path):
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception:
            with open(path, 'r', encoding='latin-1', errors='replace') as f:
                return f.read()

    def list_corpus_texts(self, exclude_path=None):
        files = []
        for name in os.listdir(CORPUS_DIR):
            if name.lower().endswith('.txt'):
                full = os.path.join(CORPUS_DIR, name)
                if exclude_path and os.path.abspath(full) == os.path.abspath(exclude_path):
                    continue
                files.append(full)
        return files

    def analyze_similarity(self, text, current_txt_path, label):
        logger.info("Folder-based similarity analysis start")
        corpus_files = self.list_corpus_texts(exclude_path=current_txt_path)
        logger.info(f"Corpus files (excluding current): {len(corpus_files)}")
        main = Text(text, label, "")
        detailed = []
        max_sim = 0.0
        total = 0
        for path in corpus_files:
            try:
                src = self._read_text_file(path)
                src_label = os.path.basename(path)
                src_text = Text(src, src_label, path)
                matcher = Matcher(main, src_text,
                                  threshold=THRESHOLD,
                                  cutoff=CUTOFF,
                                  ngramSize=NGRAM,
                                  removeStopwords=True,
                                  minDistance=MIN_DISTANCE,
                                  silent=True)
                n, infos, sim = matcher.match()
                logger.debug(f"Compared with {src_label}: matches={n}, sim={sim}")
                if sim > 0:
                    detailed.extend(infos)
                    max_sim = max(max_sim, sim)
                total += 1
            except Exception as e:
                logger.error(f"Compare error {path}: {e}")
                continue
        return {
            'similarity_score': round(max_sim, 1),
            'total_comparisons': total,
            'detailed_results': detailed[:10]
        }

    def calculate_plagiarism_score(self, similarity, ai):
        p = similarity * 0.6 + ai * 0.4
        level = 'low' if p < 15 else 'moderate' if p < 35 else 'high' if p < 60 else 'critical'
        factors = []
        if similarity > 25:
            factors.append('high_similarity_matches')
        if ai > 25:
            factors.append('ai_generated_content')
        if similarity > 15 and ai > 15:
            factors.append('multiple_detection_methods')
        if not factors:
            factors.append('low_risk_indicators')
        return {
            'plagiarism_score': round(p, 1),
            'risk_level': level,
            'contributing_factors': factors
        }

    def generate_html_report(self, student, sim, ai, plag, content):
        ts = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        def cls(score):
            return 'low-risk' if score < 15 else 'medium-risk' if score < 35 else 'high-risk' if score < 60 else 'critical-risk'
        return f"""<!DOCTYPE html>
<html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width, initial-scale=1'>
<title>Plagiarism Analysis Report - {student['student_name']}</title>
<style>
body{{font-family:'Segoe UI',Tahoma,Verdana,sans-serif;margin:0;padding:20px;background:#f8f9fa;color:#333;line-height:1.6}}
.container{{max-width:1200px;margin:0 auto;background:#fff;border-radius:10px;box-shadow:0 4px 6px rgba(0,0,0,.1);overflow:hidden}}
.header{{background:linear-gradient(135deg,#667eea,#764ba2);color:#fff;padding:30px;text-align:center}}
.header h1{{margin:0;font-size:28px;font-weight:300}}
.content{{padding:30px}}
.info-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(250px,1fr));gap:20px;margin-bottom:30px}}
.info-card{{background:#f8f9fa;padding:20px;border-radius:8px;border-left:4px solid #667eea}}
.info-card h3{{margin:0 0 10px;color:#667eea;font-size:14px;font-weight:600;text-transform:uppercase;letter-spacing:1px}}
.info-card p{{margin:0;font-size:16px;font-weight:500}}
.scores-section{{display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:30px;margin:30px 0}}
.score-card{{background:#fff;border-radius:10px;padding:25px;box-shadow:0 2px 10px rgba(0,0,0,.1);border-top:4px solid}}
.score-card.similarity{{border-top-color:#3498db}}
.score-card.ai-detection{{border-top-color:#e74c3c}}
.score-card.plagiarism{{border-top-color:#f39c12}}
.score-card h3{{margin:0 0 15px;font-size:18px;font-weight:600}}
.score-display{{font-size:36px;font-weight:700;margin:10px 0}}
.low-risk{{color:#27ae60}}.medium-risk{{color:#f39c12}}.high-risk{{color:#e74c3c}}.critical-risk{{color:#c0392b}}
.details-card{{background:#f8f9fa;border-radius:8px;padding:20px;margin-bottom:20px}}
.matches-list{{list-style:none;padding:0}}
.matches-list li{{background:#fff;padding:15px;margin-bottom:10px;border-radius:5px;border-left:3px solid #3498db}}
.risk-factors{{display:flex;flex-wrap:wrap;gap:10px;margin-top:10px}}
.risk-factor{{background:#e74c3c;color:#fff;padding:5px 12px;border-radius:20px;font-size:12px;font-weight:500}}
.footer{{background:#2c3e50;color:#fff;text-align:center;padding:20px;font-size:12px}}
.interpretation{{font-style:italic;color:#666;margin-top:5px}}
.processing-info{{background:#e8f5e8;border:1px solid #c3e6c3;border-radius:5px;padding:15px;margin-bottom:20px}}
</style></head>
<body><div class='container'><div class='header'><h1>Plagiarism Analysis Report</h1><p>Comprehensive academic integrity assessment with PDF/OCR processing</p></div>
<div class='content'>
<div class='processing-info'><h4>Document Processing Information</h4><p>This document was automatically processed using advanced PDF extraction and OCR technology. Text content includes both direct text extraction and OCR results from images within the document.</p><p><strong>Content Length:</strong> {len(content)} characters | <strong>Word Count:</strong> {len(content.split())} words</p></div>
<div class='info-grid'>
<div class='info-card'><h3>Student</h3><p>{student['student_name']}</p></div>
<div class='info-card'><h3>Classroom</h3><p>{student['classroom_name']}</p></div>
<div class='info-card'><h3>Assignment ID</h3><p>{student['assignment_id']}</p></div>
<div class='info-card'><h3>Analysis Date</h3><p>{ts}</p></div>
</div>
<div class='scores-section'>
<div class='score-card similarity'><h3>Similarity Analysis</h3><div class='score-display {cls(sim['similarity_score'])}'>{sim['similarity_score']}%</div><p>Compared against {sim['total_comparisons']} documents</p></div>
<div class='score-card ai-detection'><h3>AI Detection</h3><div class='score-display {cls(ai['ai_percentage'])}'>{ai['ai_percentage']}%</div><p class='interpretation'>{ai['interpretation']}</p><p>Confidence: {ai['confidence']}</p></div>
<div class='score-card plagiarism'><h3>Overall Plagiarism Score</h3><div class='score-display {cls(plag['plagiarism_score'])}'>{plag['plagiarism_score']}%</div><p>Risk Level: <strong>{plag['risk_level'].upper()}</strong></p><div class='risk-factors'>{' '.join([f"<span class='risk-factor'>{f.replace('_',' ').title()}</span>" for f in plag['contributing_factors']])}</div></div>
</div>
<div class='analysis-details'>
{('''<div class=\"details-card\"><h4>Similarity Matches Found</h4><ul class=\"matches-list\">''' + ''.join([f"<li><strong>{m['source']}</strong><br>Similarity: {m['similarity']}%<br><em>{m['matched_text']}</em></li>" for m in sim['detailed_results'][:5]]) + '</ul></div>') if sim['detailed_results'] else '<div class=\"details-card\"><h4>No Significant Similarity Matches Found</h4><p>The document appears to have original content with no substantial matches in our database.</p></div>'}
<div class='details-card'><h4>AI Detection Analysis</h4><p>Analyzed {ai['chunks_analyzed']} text segments with {ai['confidence']} confidence level.</p><p><strong>Interpretation:</strong> {ai['interpretation']}</p></div>
</div>
</div>
<div class='footer'><p>Generated by Enhanced Academic Integrity Analysis System with PDF/OCR Processing | {ts}</p><p>This report is confidential and intended solely for academic evaluation purposes.</p></div>
</div></body></html>"""

    def log_user_score(self, student_id, scores):
        global user_scores_log
        user_scores_log.setdefault(student_id, []).append({
            'timestamp': datetime.now().isoformat(),
            'scores': scores
        })
        user_scores_log[student_id] = user_scores_log[student_id][-10:]

    def analyze_document(self, payload):
        try:
            logger.info(f"Analyze: student_id={payload.get('student_id')} assignment_id={payload.get('assignment_id')} url={payload.get('file_url')}")
            t0 = time.time()
            content, filename = self.download_and_process_file(payload['file_url'])
            logger.info(f"Downloaded+extracted in {time.time()-t0:.2f}s | filename={filename}")
            if not content or len(content.strip()) < 10:
                raise ValueError('Could not extract meaningful content from the document')
            saved_txt_path = self.save_text_to_corpus(content, filename, payload.get('student_id'), payload.get('assignment_id'))
            sim = self.analyze_similarity(content, saved_txt_path, filename or payload['student_name'])
            ai = self.ai_detector.analyze_text(content)
            plag = self.calculate_plagiarism_score(sim['similarity_score'], ai['ai_percentage'])
            report_html = self.generate_html_report(payload, sim, ai, plag, content)
            scores = {
                'similarity_score': sim['similarity_score'],
                'ai_percentage': ai['ai_percentage'],
                'plagiarism_score': plag['plagiarism_score'],
                'content_length': len(content),
                'word_count': len(content.split())
            }
            self.log_user_score(payload['student_id'], scores)
            logger.info(f"Done | sim={scores['similarity_score']} ai={scores['ai_percentage']} plag={scores['plagiarism_score']} | corpus={sim['total_comparisons']}")
            return {
                'status': 'completed',
                'similarity_analysis': sim,
                'ai_analysis': ai,
                'plagiarism_analysis': plag,
                'report_html': report_html,
                'processing_info': {
                    'filename': filename,
                    'content_length': len(content),
                    'word_count': len(content.split()),
                    'saved_txt_path': saved_txt_path,
                    'corpus_dir': CORPUS_DIR,
                    'corpus_files_compared': sim['total_comparisons']
                },
                'errors': []
            }
        except Exception as e:
            logger.error(f"analyze_document error: {e}")
            return {
                'status': 'failed',
                'similarity_analysis': {'similarity_score': 0, 'total_comparisons': 0, 'detailed_results': []},
                'ai_analysis': {'ai_percentage': 0, 'interpretation': 'Analysis failed', 'chunks_analyzed': 0, 'confidence': 0},
                'plagiarism_analysis': {'plagiarism_score': 0, 'risk_level': 'unknown', 'contributing_factors': ['analysis_error']},
                'report_html': f"<html><body><h1>Analysis Failed</h1><p>Error: {str(e)}</p></body></html>",
                'processing_info': {'filename': 'unknown', 'content_length': 0, 'word_count': 0},
                'errors': [str(e)]
            }

api = PlagiarismAPI()

@app.route('/grade/analyze', methods=['POST'])
def analyze_submission():
    try:
        if not request.is_json:
            return jsonify({'error': 'Content-Type must be application/json'}), 400
        payload = request.get_json()
        for f in ['student_id', 'student_name', 'file_url', 'assignment_id', 'classroom_name']:
            if f not in payload:
                return jsonify({'error': f'Missing required field: {f}'}), 400
        result = api.analyze_document(payload)
        return jsonify(result)
    except Exception as e:
        logger.error(f"API error: {e}")
        return jsonify({
            'status': 'failed',
            'similarity_analysis': {'similarity_score': 0, 'total_comparisons': 0, 'detailed_results': []},
            'ai_analysis': {'ai_percentage': 0, 'interpretation': 'Analysis failed', 'chunks_analyzed': 0, 'confidence': 0},
            'plagiarism_analysis': {'plagiarism_score': 0, 'risk_level': 'unknown', 'contributing_factors': ['api_error']},
            'report_html': f"<html><body><h1>Analysis Failed</h1><p>API Error: {str(e)}</p></body></html>",
            'processing_info': {'filename': 'unknown', 'content_length': 0, 'word_count': 0},
            'errors': [str(e)]
        }), 500

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.now().isoformat(),
        'service': 'Enhanced Plagiarism Detection API with PDF/OCR Processing',
        'capabilities': [
            'PDF text extraction',
            'OCR for image-based documents',
            'DOCX/PPTX processing',
            'Image preprocessing',
            'Multi-format support'
        ]
    })

@app.route('/test-processing', methods=['POST'])
def test_file_processing():
    try:
        if not request.is_json:
            return jsonify({'error': 'Content-Type must be application/json'}), 400
        payload = request.get_json()
        if 'file_url' not in payload:
            return jsonify({'error': 'Missing required field: file_url'}), 400
        logger.info(f"Test-processing URL: {payload.get('file_url')}")
        content, filename = api.download_and_process_file(payload['file_url'])
        return jsonify({
            'status': 'success',
            'filename': filename,
            'content_preview': content[:500] + '...' if len(content) > 500 else content,
            'content_length': len(content),
            'word_count': len(content.split())
        })
    except Exception as e:
        logger.error(f"Test processing error: {e}")
        return jsonify({'status': 'failed', 'error': str(e)}), 500

@app.route('/user/<student_id>/scores', methods=['GET'])
def get_user_scores(student_id):
    if student_id in user_scores_log:
        return jsonify({'student_id': student_id, 'scores_history': user_scores_log[student_id]})
    else:
        return jsonify({'student_id': student_id, 'scores_history': [], 'message': 'No scores found for this user'})

@app.route('/stats', methods=['GET'])
def get_stats():
    total_users = len(user_scores_log)
    total_analyses = sum(len(scores) for scores in user_scores_log.values())
    try:
        corpus_count = len([n for n in os.listdir(CORPUS_DIR) if n.lower().endswith('.txt')])
    except Exception:
        corpus_count = 0
    return jsonify({
        'total_users_analyzed': total_users,
        'total_analyses_performed': total_analyses,
        'corpus_text_files': corpus_count,
        'system_status': 'operational',
        'features': [
            'PDF text extraction with PyMuPDF',
            'OCR processing with Tesseract',
            'Image preprocessing with OpenCV',
            'DOCX/PPTX support',
            'Multi-format document processing',
            'Persistent text corpus and folder-based similarity'
        ]
    })

if __name__ == '__main__':
    os.makedirs(LOGS_DIR, exist_ok=True)
    os.makedirs(CORPUS_DIR, exist_ok=True)
    print('Starting Enhanced Plagiarism Detection API Service with PDF/OCR Processing...')
    print('Server will run on http://localhost:5000')
    print('\nFeatures enabled:')
    print('✓ PDF text extraction with PyMuPDF')
    print('✓ OCR processing with Tesseract')
    print('✓ Image preprocessing with OpenCV/PIL')
    print('✓ DOCX/PPTX document support')
    print('✓ Multi-format file processing')
    print('\nAvailable endpoints:')
    print('POST /grade/analyze - Main plagiarism analysis endpoint')
    print('POST /test-processing - Test file processing capabilities')
    print('GET /health - Health check with capabilities')
    print('GET /user/<student_id>/scores - Get user\'s score history')
    print('GET /stats - Get system statistics')
    try:
        pytesseract.get_tesseract_version()
        print(f"✓ Tesseract OCR ready: {pytesseract.get_tesseract_version()}")
    except Exception as e:
        print(f"⚠ Tesseract warning: {e}")
        print('  OCR functionality may be limited')
    app.run(host='0.0.0.0', port=5000, debug=True)
import mimetypes
from werkzeug.utils import secure_filename

# PDF/OCR Processing imports
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import fitz  # PyMuPDF for handling PDFs
import cv2
import numpy as np
import io
import traceback
from docx import Document
from pptx import Presentation
import csv

# Initialize Flask app
app = Flask(__name__)

# -------------------------------------------------------------
# Global configuration (easy to tweak)
# -------------------------------------------------------------
THRESHOLD = 3       # --threshold
CUTOFF = 5          # --cutoff
NGRAM = 3           # --ngram
MIN_DISTANCE = 8    # --distance
CORPUS_DIR = os.path.join(os.path.dirname(__file__), 'corpus')
LOGS_DIR = os.path.join(os.path.dirname(__file__), 'logs')
LOG_FILE = os.path.join(LOGS_DIR, 'api.log')

# Ensure base folders exist
os.makedirs(CORPUS_DIR, exist_ok=True)
os.makedirs(LOGS_DIR, exist_ok=True)

# Configure logging (console + rotating file)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger('plagiarism_api')
logger.setLevel(logging.DEBUG)
if not logger.handlers:
    file_handler = RotatingFileHandler(LOG_FILE, maxBytes=2_000_000, backupCount=5, encoding='utf-8')
    file_handler.setLevel(logging.DEBUG)
    file_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(name)s - %(message)s'))
    logger.addHandler(file_handler)
    # Mirror to console as well
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)
    console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
    logger.addHandler(console_handler)

# Global storage for user scores
            try:
                nltk.data.find('tokenizers/punkt')
            except LookupError:
                nltk.download('punkt')
            try:
                nltk.data.find('corpora/stopwords')
            except LookupError:
                nltk.download('stopwords')

        def download_and_process_file(self, url):
            """Download file from URL and extract text content"""
            try:
                response = requests.get(url, timeout=60)
                response.raise_for_status()

                # Create temporary file
                temp_dir = tempfile.mkdtemp()

                # Determine filename from URL or Content-Disposition header
                filename = None
                if 'content-disposition' in response.headers:
                    cd = response.headers['content-disposition']
                    filename_match = re.findall('filename=(.+)', cd)
                    if filename_match:
                        filename = filename_match[0].strip('"')

                if not filename:
                    filename = os.path.basename(urlparse(url).path) or "document"

                # Ensure filename has proper extension based on content type
                content_type = response.headers.get('content-type', '').lower()
                if not os.path.splitext(filename)[1]:
                    if 'pdf' in content_type:
                        filename += '.pdf'
                    elif 'word' in content_type or 'docx' in content_type:
                        filename += '.docx'
                    elif 'powerpoint' in content_type or 'pptx' in content_type:
                        filename += '.pptx'
                    else:
                        filename += '.txt'

                temp_path = os.path.join(temp_dir, secure_filename(filename))

                logger.info(f"Downloading file from URL: {url}")
                logger.debug(f"Response headers: {dict(response.headers)}")
                # Save file content
                with open(temp_path, 'wb') as f:
                    f.write(response.content)
                logger.info(f"Saved download to temp path: {temp_path}")

                # Extract text using PDF processor
                start_extract = time.time()
                extracted_text = self.pdf_processor.process_file_content(temp_path, content_type)
                elapsed_extract = time.time() - start_extract
                logger.info(f"Extraction complete in {elapsed_extract:.2f}s, extracted length={len(extracted_text)} chars")

                # Clean up temporary files
                try:
                    os.unlink(temp_path)
                    os.rmdir(temp_dir)
                except Exception:
                    pass

                return extracted_text, filename

            except Exception as e:
                logger.error(f"Error downloading and processing file from {url}: {e}")
                raise
        
        except Exception as e:
            logger.error(f"Error preprocessing image: {e}")
            # If preprocessing fails, return the original image
            return image, image
    
    def perform_ocr(self, image, lang="eng", config='--psm 6'):
        """
        Perform OCR on the given image with error handling.
        Returns empty string if OCR fails.
        """
        try:
            # Try to improve OCR with image preprocessing
            _, binary_image = self.preprocess_image(image)
            
            # Perform OCR with confidence
            ocr_data = pytesseract.image_to_data(
                binary_image, lang=lang, config=config, output_type=pytesseract.Output.DICT
            )
            
            # Filter out low confidence results
            text_parts = []
            for i, conf in enumerate(ocr_data['conf']):
                if float(conf) > 30:  # Only include text with confidence above 30%
                    text = ocr_data['text'][i]
                    if text.strip():
                        text_parts.append(text)
            
            # Join text parts
            text = " ".join(text_parts)
            
            # Clean up text
            text = re.sub(r'\s+', ' ', text).strip()
            
            return text
        except Exception as e:
            logger.error(f"OCR error: {e}")
            return ""
    
    def extract_text_from_pdf(self, file_path):
        """
        Extract text and images from PDF using PyMuPDF with OCR for images.
        """
        try:
            # Open PDF with PyMuPDF
            pdf_document = fitz.open(file_path)
            extracted_text = ""
            
            # Process each page
            for page_index in range(len(pdf_document)):
                page = pdf_document.load_page(page_index)
                
                # Extract direct text first
                page_text = page.get_text()
                
                # If page has substantial text, use it
                if len(page_text.strip()) > 50:
                    extracted_text += f"\n\n--- Page {page_index + 1} ---\n\n"
                    extracted_text += page_text
                else:
                    # Page might be image-based, try OCR on the entire page
                    try:
                        # Render page as image
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        
                        # Perform OCR on the page
                        ocr_text = self.perform_ocr(img)
                        
                        if ocr_text.strip():
                            extracted_text += f"\n\n--- Page {page_index + 1} (OCR) ---\n\n"
                            extracted_text += ocr_text
                        
                    except Exception as e:
                        logging.error(f"Error processing page {page_index} with OCR: {e}")
                
                # Extract images and perform OCR on them
                image_list = page.get_images(full=True)
                
                for img_index, img_info in enumerate(image_list):
                    try:
                        xref = img_info[0]
                        base_image = fitz.Pixmap(pdf_document, xref)
                        
                        # Convert to PIL Image
                        if base_image.n < 5:  # RGB
                            pil_img = Image.frombytes("RGB", [base_image.width, base_image.height], base_image.samples)
                        else:  # CMYK
                            cmyk_image = fitz.Pixmap(fitz.csRGB, base_image)
                            pil_img = Image.frombytes("RGB", [cmyk_image.width, cmyk_image.height], cmyk_image.samples)
                            cmyk_image = None  # Free memory
                        
                        # Skip small images (likely icons or decorations)
                        if pil_img.width < 50 or pil_img.height < 50:
                            continue
                        
                        # Extract text from image using OCR
                        ocr_text = self.perform_ocr(pil_img)
                        
                        if ocr_text.strip():
                            extracted_text += f"\n\n--- Page {page_index + 1}, Image {img_index + 1} (OCR) ---\n\n"
                            extracted_text += ocr_text
                        
                    except Exception as e:
                        logging.error(f"Error processing image {img_index} on page {page_index}: {e}")
                        continue
            
            pdf_document.close()
            return extracted_text.strip()
            
        except Exception as e:
            logger.error(f"Error extracting from PDF {file_path}: {str(e)}\n{traceback.format_exc()}")
            return f"Error extracting content from PDF: {str(e)}"
    
    def extract_text_from_docx(self, file_path):
        """Extract text from DOCX files."""
        try:
            doc = Document(file_path)
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error extracting from DOCX {file_path}: {e}")
            return f"Error extracting content from DOCX: {str(e)}"
    
    def extract_text_from_pptx(self, file_path):
        """Extract text from PPTX files."""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {slide_num + 1} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error extracting from PPTX {file_path}: {e}")
            return f"Error extracting content from PPTX: {str(e)}"
    
    def extract_text_from_txt(self, file_path):
        """Extract text from plain text files with encoding detection."""
        try:
            # Try multiple encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try binary mode with error handling
            with open(file_path, "rb") as f:
                binary = f.read()
                return binary.decode('utf-8', errors='replace')
                
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return f"Error reading text file: {str(e)}"
    
    def extract_text_from_csv(self, file_path):
        """Extract text from CSV files."""
        try:
            text_content = []
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as csv_file:
                        reader = csv.reader(csv_file)
                        row_count = 0
                        
                        for row in reader:
                            row_count += 1
                            if row_count > 100:  # Limit to 100 rows
                                text_content.append("...truncated, showing first 100 rows...")
                                break
                            text_content.append(" ".join(row))
                    
                    return "\n".join(text_content)
                    
                except UnicodeDecodeError:
                    continue
            
            return "Error: Could not decode CSV file"
            
        except Exception as e:
            logger.error(f"Error reading CSV file {file_path}: {e}")
            return f"Error reading CSV file: {str(e)}"
    
    def extract_text_from_image(self, file_path):
        """Extract text from image files using OCR."""
        try:
            img = Image.open(file_path)
            ocr_text = self.perform_ocr(img)
            return ocr_text if ocr_text.strip() else "No text detected in image"
            
        except Exception as e:
            logger.error(f"Error processing image file {file_path}: {e}")
            return f"Error processing image file: {str(e)}"
    
    def process_file_content(self, file_path, content_type=None):
        """
        Process a file and extract text content based on file type.
        """
        try:
            # Determine file type from extension if content_type not provided
            if not content_type:
                _, ext = os.path.splitext(file_path)
                ext = ext.lower().lstrip('.')
            else:
                # Extract extension from content type or file path
                if '/' in content_type:
                    ext = content_type.split('/')[-1]
                else:
                    _, ext = os.path.splitext(file_path)
                    ext = ext.lower().lstrip('.')
            
            # Route to appropriate extraction method
            if ext == 'pdf':
                return self.extract_text_from_pdf(file_path)
            elif ext in ['docx', 'doc']:
                return self.extract_text_from_docx(file_path)
            elif ext in ['pptx', 'ppt']:
                return self.extract_text_from_pptx(file_path)
            elif ext == 'txt':
                return self.extract_text_from_txt(file_path)
            elif ext == 'csv':
                return self.extract_text_from_csv(file_path)
            elif ext in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif']:
                return self.extract_text_from_image(file_path)
            else:
                # Try to read as text file for unknown types
                logging.warning(f"Unknown file type: {ext}. Attempting text extraction.")
                return self.extract_text_from_txt(file_path)
                
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return f"Error processing file: {str(e)}"


class Text:
    def __init__(self, raw_text, label, filepath, removeStopwords=True):
        if isinstance(raw_text, list):
            self.text = ' \n '.join(raw_text)
        else:
            self.text = raw_text
        self.label = label
        self.filepath = filepath
        self.preprocess(self.text)
        self.tokens = self.getTokens(removeStopwords)
        self.trigrams = self.ngrams(3)
        self.checksum = self.calculate_checksum()

    def calculate_checksum(self):
        """Calculate a checksum for the file content to detect changes"""
        return hashlib.md5(self.text.encode()).hexdigest()

    def preprocess(self, text):
        """ Heals hyphenated words, and maybe other things. """
        self.text = re.sub(r'([A-Za-z])- ([a-z])', r'\1\2', text)

    def getTokens(self, removeStopwords=True):
        """ Tokenizes the text, breaking it up into words, removing punctuation. """
        tokenizer = nltk.RegexpTokenizer('[a-zA-Z]\\w+\'?\\w*')
        spans = list(tokenizer.span_tokenize(self.text))
        
        if spans:
            self.length = spans[-1][-1]
        else:
            self.length = 0
            
        tokens = tokenizer.tokenize(self.text)
        tokens = [token.lower() for token in tokens]
        stemmer = LancasterStemmer()
        tokens = [stemmer.stem(token) for token in tokens]
        
        if not removeStopwords:
            self.spans = spans
            return tokens
            
        tokenSpans = list(zip(tokens, spans))
        stopwords = nltk.corpus.stopwords.words('english')
        tokenSpans = [token for token in tokenSpans if token[0] not in stopwords]
        self.spans = [x[1] for x in tokenSpans]
        return [x[0] for x in tokenSpans]

    def ngrams(self, n):
        """ Returns ngrams for the text."""
        return list(ngrams(self.tokens, n))


class ExtendedMatch:
    def __init__(self, a, b, sizeA, sizeB):
        self.a = a
        self.b = b
        self.sizeA = sizeA
        self.sizeB = sizeB
        self.healed = False
        self.extendedBackwards = 0
        self.extendedForwards = 0

    def __repr__(self):
        out = "a: %s, b: %s, size a: %s, size b: %s" % (self.a, self.b, self.sizeA, self.sizeB)
        if self.extendedBackwards:
            out += ", extended backwards x%s" % self.extendedBackwards
        if self.extendedForwards:
            out += ", extended forwards x%s" % self.extendedForwards
        if self.healed:
            out += ", healed"
        return out


class Matcher:
    def __init__(self, textObjA, textObjB, threshold=3, cutoff=5, ngramSize=3, removeStopwords=True, minDistance=8, silent=True):
        self.threshold = threshold
        self.ngramSize = ngramSize
        self.minDistance = minDistance
        self.silent = silent

        self.textA = textObjA
        self.textB = textObjB

        self.textAgrams = self.textA.ngrams(ngramSize)
        self.textBgrams = self.textB.ngrams(ngramSize)

        self.locationsA = []
        self.locationsB = []
        self.match_texts = []

        self.initial_matches = self.get_initial_matches()
        self.healed_matches = self.heal_neighboring_matches()
        self.extended_matches = self.extend_matches()
        self.extended_matches = [match for match in self.extended_matches
                                if min(match.sizeA, match.sizeB) >= cutoff]

        self.numMatches = len(self.extended_matches)
        self.similarity_score = self.calculate_similarity()

    def calculate_similarity(self):
        if not self.extended_matches:
            return 0.0
            
        total_matched_tokens_A = sum(match.sizeA for match in self.extended_matches)
        total_tokens_A = len(self.textA.tokens)
        
        if total_tokens_A == 0:
            return 0.0
            
        similarity = (total_matched_tokens_A / total_tokens_A) * 100
        return round(similarity, 2)

    def get_initial_matches(self):
        sequence = SequenceMatcher(None, self.textAgrams, self.textBgrams)
        matchingBlocks = sequence.get_matching_blocks()
        highMatchingBlocks = [match for match in matchingBlocks if match.size > self.threshold]
        return highMatchingBlocks

    def getTokensText(self, text, start, length):
        if start < 0:
            start = 0
        
        matchTokens = text.tokens[start:start + length]
        
        if start >= len(text.spans) or start + length > len(text.spans):
            return ""
            
        spans = text.spans[start:start + length]
        if len(spans) == 0:
            passage = ""
        else:
            passage = text.text[spans[0][0]:spans[-1][-1]]
        return passage

    def getLocations(self, text, start, length, asPercentages=False):
        if start >= len(text.spans) or start + length > len(text.spans):
            return None
            
        spans = text.spans[start:start + length]
        if len(spans) == 0:
            return None
            
        if asPercentages:
            locations = (spans[0][0] / text.length, spans[-1][-1] / text.length)
        else:
            try:
                locations = (spans[0][0], spans[-1][-1])
            except IndexError:
                return None
        return locations

    def heal_neighboring_matches(self):
        healedMatches = []
        ignoreNext = False
        matches = self.initial_matches.copy()
        
        if len(matches) == 1:
            match = matches[0]
            sizeA, sizeB = match.size, match.size
            match = ExtendedMatch(match.a, match.b, sizeA, sizeB)
            healedMatches.append(match)
            return healedMatches
            
        for i, match in enumerate(matches):
            if i + 1 > len(matches) - 1:
                break
            nextMatch = matches[i + 1]
            
            if ignoreNext:
                ignoreNext = False
                continue
            else:
                if (nextMatch.a - (match.a + match.size)) < self.minDistance:
                    sizeA = (nextMatch.a + nextMatch.size) - match.a
                    sizeB = (nextMatch.b + nextMatch.size) - match.b
                    healed = ExtendedMatch(match.a, match.b, sizeA, sizeB)
                    healed.healed = True
                    healedMatches.append(healed)
                    ignoreNext = True
                else:
                    sizeA, sizeB = match.size, match.size
                    match = ExtendedMatch(match.a, match.b, sizeA, sizeB)
                    healedMatches.append(match)
        return healedMatches

    def edit_ratio(self, wordA, wordB):
        distance = editDistance(wordA, wordB)
        averageLength = (len(wordA) + len(wordB)) / 2
        return distance / averageLength

    def extend_matches(self, cutoff=0.4):
        extended = False
        for match in self.healed_matches:
            if match.a > 0 and match.b > 0 and len(self.textAgrams) > match.a - 1 and len(self.textBgrams) > match.b - 1:
                wordA = self.textAgrams[(match.a - 1)][0]
                wordB = self.textBgrams[(match.b - 1)][0]
                if self.edit_ratio(wordA, wordB) < cutoff:
                    match.a -= 1
                    match.b -= 1
                    match.sizeA += 1
                    match.sizeB += 1
                    match.extendedBackwards += 1
                    extended = True
                    
            idxA = match.a + match.sizeA + 1
            idxB = match.b + match.sizeB + 1
            if idxA >= len(self.textAgrams) or idxB >= len(self.textBgrams):
                continue
                
            wordA = self.textAgrams[idxA][-1] if idxA < len(self.textAgrams) else ""
            wordB = self.textBgrams[idxB][-1] if idxB < len(self.textBgrams) else ""
            
            if wordA and wordB and self.edit_ratio(wordA, wordB) < cutoff:
                match.sizeA += 1
                match.sizeB += 1
                match.extendedForwards += 1
                extended = True

        if extended:
            self.extend_matches()

        return self.healed_matches

    def match(self):
        matches_info = []
        for num, match in enumerate(self.extended_matches):
            lengthA = match.sizeA + self.ngramSize - 1
            matched_text = self.getTokensText(self.textA, match.a, lengthA)
            
            if matched_text:
                matches_info.append({
                    "source": self.textB.label,
                    "similarity": self.similarity_score,
                    "matched_text": matched_text[:200] + "..." if len(matched_text) > 200 else matched_text
                })

        return self.numMatches, matches_info, self.similarity_score


class AIDetector:
    """Mock AI detection class - generates random but realistic AI scores"""
    
    def __init__(self):
        self.patterns = [
            "repetitive sentence structures",
            "uniform vocabulary complexity",
            "lack of personal voice",
            "perfect grammar consistency",
            "unnatural transitions",
            "generic examples"
        ]
    
    def analyze_text(self, text):
        """Analyze text for AI-generated content"""
        # Mock analysis - in reality, this would use AI detection models
        word_count = len(text.split())
        
        # Generate realistic AI percentage based on text characteristics
        base_score = random.uniform(5.0, 35.0)
        
        # Adjust based on text length (longer texts tend to have lower AI scores)
        if word_count > 1000:
            base_score *= 0.8
        elif word_count < 200:
            base_score *= 1.2
            
        ai_percentage = min(95.0, max(0.0, base_score))
        
        # Generate interpretation
        if ai_percentage < 15:
            interpretation = "Low AI probability"
        elif ai_percentage < 35:
            interpretation = "Moderate AI probability"
        elif ai_percentage < 60:
            interpretation = "High AI probability"
        else:
            interpretation = "Very high AI probability"
            
        # Mock chunks analyzed
        chunks_analyzed = max(1, word_count // 50)
        
        # Mock confidence (higher confidence for more extreme scores)
        if ai_percentage < 20 or ai_percentage > 70:
            confidence = random.uniform(0.8, 0.95)
        else:
            confidence = random.uniform(0.6, 0.85)
            
        return {
            "ai_percentage": round(ai_percentage, 1),
            "interpretation": interpretation,
            "chunks_analyzed": chunks_analyzed,
            "confidence": round(confidence, 2)
        }


class PlagiarismAPI:
    def __init__(self):
        self.ai_detector = AIDetector()
        self.pdf_processor = PDFProcessor()
        self.ensure_nltk_data()
        
    def ensure_nltk_data(self):
        """Download required NLTK data"""
    try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords')
    try:
    def download_and_process_file(self, url):
        """Download file from URL and extract text content"""
        try:
            response = requests.get(url, timeout=60)
            response.raise_for_status()
            
            # Create temporary file
            temp_dir = tempfile.mkdtemp()
            
        cd = response.headers['content-disposition']
            filename = None
            if 'content-disposition' in response.headers:
                import re
                cd = response.headers['content-disposition']
                filename_match = re.findall('filename=(.+)', cd)
                if filename_match:
                    filename = filename_match[0].strip('"')
            
            if not filename:
                filename = os.path.basename(urlparse(url).path) or "document"
                
            # Ensure filename has proper extension based on content type
            content_type = response.headers.get('content-type', '').lower()
            if not os.path.splitext(filename)[1]:
                if 'pdf' in content_type:
                    filename += '.pdf'
                elif 'word' in content_type or 'docx' in content_type:
                    filename += '.docx'
                elif 'powerpoint' in content_type or 'pptx' in content_type:
                    filename += '.pptx'
                else:
                    filename += '.txt'
            
            temp_path = os.path.join(temp_dir, secure_filename(filename))
            
            logger.info(f"Downloading file from URL: {url}")
            logger.debug(f"Response headers: {dict(response.headers)}")
            # Save file content
            with open(temp_path, 'wb') as f:
                f.write(response.content)
            logger.info(f"Saved download to temp path: {temp_path}")

            # Extract text using PDF processor
            start_extract = time.time()
            extracted_text = self.pdf_processor.process_file_content(temp_path, content_type)
            elapsed_extract = time.time() - start_extract
            logger.info(f"Extraction complete in {elapsed_extract:.2f}s, extracted length={len(extracted_text)} chars")
            
            # Clean up temporary files
            try:
                os.unlink(temp_path)
                os.rmdir(temp_dir)
            except:
                pass
            
            return extracted_text, filename
            
        except Exception as e:
            logger.error(f"Error downloading and processing file from {url}: {e}")
            raise
    
    def save_text_to_corpus(self, text_content, original_filename, student_id, assignment_id):
        """Save extracted text to the persistent corpus directory as a .txt file and return its path"""
        # Sanitize name
        base_name = os.path.splitext(original_filename or 'document')[0]
        base_name = re.sub(r'[^A-Za-z0-9._-]+', '_', base_name)[:80]
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        checksum = hashlib.md5(text_content.encode('utf-8', errors='ignore')).hexdigest()[:8]
        student_id_safe = re.sub(r'[^A-Za-z0-9_-]+', '_', str(student_id))
        assignment_id_safe = re.sub(r'[^A-Za-z0-9_-]+', '_', str(assignment_id))
        filename = f"{base_name}__{student_id_safe}__{assignment_id_safe}__{timestamp}__{checksum}.txt"
        txt_path = os.path.join(CORPUS_DIR, filename)

        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(text_content)

        logger.info(f"Saved extracted text to corpus: {txt_path}")
        return txt_path

    def _read_text_file(self, path):
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception:
            # Fallback encodings
            with open(path, 'r', encoding='latin-1', errors='replace') as f:
                return f.read()

    def list_corpus_texts(self, exclude_path=None):
        """List all .txt files in the corpus directory, optionally excluding one path"""
        files = []
        for name in os.listdir(CORPUS_DIR):
            if not name.lower().endswith('.txt'):
                continue
            full = os.path.join(CORPUS_DIR, name)
            if exclude_path and os.path.abspath(full) == os.path.abspath(exclude_path):
                continue
            files.append(full)
        return files

    def analyze_similarity(self, text_content, current_txt_path, current_label):
        """Analyze text similarity against all .txt files in the corpus directory (max across sources)."""
        logger.info("Starting folder-based similarity analysis")
        corpus_files = self.list_corpus_texts(exclude_path=current_txt_path)
        logger.info(f"Corpus files found (excluding current): {len(corpus_files)}")

        main_text = Text(text_content, current_label, "")
        detailed_results = []
        max_similarity = 0.0
        total_comparisons = 0

        for path in corpus_files:
            try:
                source_content = self._read_text_file(path)
                source_name = os.path.basename(path)
                source_text = Text(source_content, source_name, path)
                matcher = Matcher(
                    main_text,
                    source_text,
                    threshold=THRESHOLD,
                    cutoff=CUTOFF,
                    ngramSize=NGRAM,
                    removeStopwords=True,
                    minDistance=MIN_DISTANCE,
                    silent=True
                )
                num_matches, matches_info, similarity = matcher.match()
                logger.debug(f"Compared with {source_name} -> matches={num_matches}, similarity={similarity}")

                if similarity > 0:
                    # Keep snippets from this source
                    detailed_results.extend(matches_info)
                    max_similarity = max(max_similarity, similarity)

                total_comparisons += 1

            except Exception as e:
                logger.error(f"Error comparing with {path}: {e}")
                continue

        # Limit to top 10 detailed results to keep payload small
        return {
            "similarity_score": round(max_similarity, 1),
            "total_comparisons": total_comparisons,
            "detailed_results": detailed_results[:10]
        }
    
    def calculate_plagiarism_score(self, similarity_score, ai_percentage):
        """Calculate overall plagiarism score"""
        # Weight the scores (you can adjust these weights)
        similarity_weight = 0.6
        ai_weight = 0.4
        
        plagiarism_score = (similarity_score * similarity_weight) + (ai_percentage * ai_weight)
        
        # Determine risk level
        if plagiarism_score < 15:
            risk_level = "low"
        elif plagiarism_score < 35:
            risk_level = "moderate"
        elif plagiarism_score < 60:
            risk_level = "high"
        else:
            risk_level = "critical"
            
        # Determine contributing factors
        contributing_factors = []
        if similarity_score > 25:
            contributing_factors.append("high_similarity_matches")
        if ai_percentage > 25:
            contributing_factors.append("ai_generated_content")
        if similarity_score > 15 and ai_percentage > 15:
            contributing_factors.append("multiple_detection_methods")
        if not contributing_factors:
            contributing_factors.append("low_risk_indicators")
            
        return {
            "plagiarism_score": round(plagiarism_score, 1),
            "risk_level": risk_level,
            "contributing_factors": contributing_factors
        }
    
    def generate_html_report(self, student_data, similarity_analysis, ai_analysis, plagiarism_analysis, content):
        """Generate comprehensive HTML report"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # Determine color classes based on scores
        def get_score_class(score):
            if score < 15:
                return "low-risk"
            elif score < 35:
                return "medium-risk"
            elif score < 60:
                return "high-risk"
            else:
                return "critical-risk"
        
        html_report = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Plagiarism Analysis Report - {student_data['student_name']}</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 20px;
            background-color: #f8f9fa;
            color: #333;
            line-height: 1.6;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            overflow: hidden;
        }}
        .header {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 30px;
            text-align: center;
        }}
        .header h1 {{
            margin: 0;
            font-size: 28px;
            font-weight: 300;
        }}
        .content {{
            padding: 30px;
        }}
        .info-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }}
        .info-card {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            border-left: 4px solid #667eea;
        }}
        .info-card h3 {{
            margin: 0 0 10px 0;
            color: #667eea;
            font-size: 14px;
            font-weight: 600;
            text-transform: uppercase;
            letter-spacing: 1px;
        }}
        .info-card p {{
            margin: 0;
            font-size: 16px;
            font-weight: 500;
        }}
        .scores-section {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 30px;
            margin: 30px 0;
        }}
        .score-card {{
            background: white;
            border-radius: 10px;
            padding: 25px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
            border-top: 4px solid;
        }}
        .score-card.similarity {{
            border-top-color: #3498db;
        }}
        .score-card.ai-detection {{
            border-top-color: #e74c3c;
        }}
        .score-card.plagiarism {{
            border-top-color: #f39c12;
        }}
        .score-card h3 {{
            margin: 0 0 15px 0;
            font-size: 18px;
            font-weight: 600;
        }}
        .score-display {{
            font-size: 36px;
            font-weight: bold;
            margin: 10px 0;
        }}
        .low-risk {{ color: #27ae60; }}
        .medium-risk {{ color: #f39c12; }}
        .high-risk {{ color: #e74c3c; }}
        .critical-risk {{ color: #c0392b; }}
        .analysis-details {{
            margin-top: 30px;
        }}
        .details-card {{
            background: #f8f9fa;
            border-radius: 8px;
            padding: 20px;
            margin-bottom: 20px;
        }}
        .details-card h4 {{
            margin: 0 0 15px 0;
            color: #2c3e50;
        }}
        .matches-list {{
            list-style: none;
            padding: 0;
        }}
        .matches-list li {{
            background: white;
            padding: 15px;
            margin-bottom: 10px;
            border-radius: 5px;
            border-left: 3px solid #3498db;
        }}
        .risk-factors {{
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            margin-top: 10px;
        }}
        .risk-factor {{
            background: #e74c3c;
            color: white;
            padding: 5px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: 500;
        }}
        .footer {{
            background: #2c3e50;
            color: white;
            text-align: center;
            padding: 20px;
            font-size: 12px;
        }}
        .interpretation {{
            font-style: italic;
            color: #666;
            margin-top: 5px;
        }}
        .processing-info {{
            background: #e8f5e8;
            border: 1px solid #c3e6c3;
            border-radius: 5px;
            padding: 15px;
            margin-bottom: 20px;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Plagiarism Analysis Report</h1>
            <p>Comprehensive academic integrity assessment with PDF/OCR processing</p>
        </div>
        
        <div class="content">
            <div class="processing-info">
                <h4>Document Processing Information</h4>
                <p>This document was automatically processed using advanced PDF extraction and OCR technology. 
                Text content includes both direct text extraction and OCR results from images within the document.</p>
                <p><strong>Content Length:</strong> {len(content)} characters | 
                <strong>Word Count:</strong> {len(content.split())} words</p>
            </div>
            
            <div class="info-grid">
                <div class="info-card">
                    <h3>Student</h3>
                    <p>{student_data['student_name']}</p>
                </div>
                <div class="info-card">
                    <h3>Classroom</h3>
                    <p>{student_data['classroom_name']}</p>
                </div>
                <div class="info-card">
                    <h3>Assignment ID</h3>
                    <p>{student_data['assignment_id']}</p>
                </div>
                <div class="info-card">
                    <h3>Analysis Date</h3>
                    <p>{timestamp}</p>
                </div>
            </div>
            
            <div class="scores-section">
                <div class="score-card similarity">
                    <h3>Similarity Analysis</h3>
                    <div class="score-display {get_score_class(similarity_analysis['similarity_score'])}">
                        {similarity_analysis['similarity_score']}%
                    </div>
                    <p>Compared against {similarity_analysis['total_comparisons']} documents</p>
                </div>
                
                <div class="score-card ai-detection">
                    <h3>AI Detection</h3>
                    <div class="score-display {get_score_class(ai_analysis['ai_percentage'])}">
                        {ai_analysis['ai_percentage']}%
                    </div>
                    <p class="interpretation">{ai_analysis['interpretation']}</p>
                    <p>Confidence: {ai_analysis['confidence']}</p>
                </div>
                
                <div class="score-card plagiarism">
                    <h3>Overall Plagiarism Score</h3>
                    <div class="score-display {get_score_class(plagiarism_analysis['plagiarism_score'])}">
                        {plagiarism_analysis['plagiarism_score']}%
                    </div>
                    <p>Risk Level: <strong>{plagiarism_analysis['risk_level'].upper()}</strong></p>
                    <div class="risk-factors">
                        {' '.join([f'<span class="risk-factor">{factor.replace("_", " ").title()}</span>' for factor in plagiarism_analysis['contributing_factors']])}
                    </div>
                </div>
            </div>
            
            <div class="analysis-details">
                {f'''
                <div class="details-card">
                    <h4>Similarity Matches Found</h4>
                    <ul class="matches-list">
                        {"".join([f"<li><strong>{match['source']}</strong><br>Similarity: {match['similarity']}%<br><em>{match['matched_text']}</em></li>" for match in similarity_analysis['detailed_results'][:5]])}
                    </ul>
                </div>
                ''' if similarity_analysis['detailed_results'] else '<div class="details-card"><h4>No Significant Similarity Matches Found</h4><p>The document appears to have original content with no substantial matches in our database.</p></div>'}
                
                <div class="details-card">
                    <h4>AI Detection Analysis</h4>
                    <p>Analyzed {ai_analysis['chunks_analyzed']} text segments with {ai_analysis['confidence']} confidence level.</p>
                    <p><strong>Interpretation:</strong> {ai_analysis['interpretation']}</p>
                </div>
            </div>
        </div>
        
        <div class="footer">
            <p>Generated by Enhanced Academic Integrity Analysis System with PDF/OCR Processing | {timestamp}</p>
            <p>This report is confidential and intended solely for academic evaluation purposes.</p>
        </div>
    </div>
</body>
</html>"""
        
        return html_report
    
    def log_user_score(self, student_id, scores):
        """Log user scores for tracking"""
    global user_scores_log
        
        if student_id not in user_scores_log:
            user_scores_log[student_id] = []
            
        user_scores_log[student_id].append({
            "timestamp": datetime.now().isoformat(),
            "scores": scores
        })
        
        # Keep only last 10 entries per user
        user_scores_log[student_id] = user_scores_log[student_id][-10:]
    
    def analyze_document(self, payload):
        """Main analysis function with integrated PDF/OCR processing"""
        try:
            # Download and process file with enhanced PDF/OCR capabilities
            logger.info(f"Analyze request: student_id={payload.get('student_id')} assignment_id={payload.get('assignment_id')} file_url={payload.get('file_url')}")
            t0 = time.time()
            content, filename = self.download_and_process_file(payload['file_url'])
            logger.info(f"Downloaded and extracted in {(time.time()-t0):.2f}s | raw filename={filename}")
            
            # Validate that we got content
            if not content or len(content.strip()) < 10:
                raise ValueError("Could not extract meaningful content from the document")

            # Persist extracted text to corpus for future comparisons
            saved_txt_path = self.save_text_to_corpus(
                content,
                original_filename=filename,
                student_id=payload.get('student_id'),
                assignment_id=payload.get('assignment_id')
            )
            
            # Perform similarity analysis against corpus folder
            similarity_analysis = self.analyze_similarity(content, saved_txt_path, filename or payload['student_name'])
            
            # Perform AI detection analysis
            ai_analysis = self.ai_detector.analyze_text(content)
            
            # Calculate overall plagiarism score
            plagiarism_analysis = self.calculate_plagiarism_score(
                similarity_analysis['similarity_score'],
                ai_analysis['ai_percentage']
            )
            
            # Generate HTML report
            html_report = self.generate_html_report(
                payload, similarity_analysis, ai_analysis, plagiarism_analysis, content
            )
            
            # Log scores
            scores = {
                "similarity_score": similarity_analysis['similarity_score'],
                "ai_percentage": ai_analysis['ai_percentage'],
                "plagiarism_score": plagiarism_analysis['plagiarism_score'],
                "content_length": len(content),
                "word_count": len(content.split())
            }
        self.log_user_score(payload['student_id'], scores)
        logger.info(f"Analysis completed | similarity={scores['similarity_score']} AI={scores['ai_percentage']} plagiarism={scores['plagiarism_score']} | corpus_count={similarity_analysis['total_comparisons']}")
            
            return {
                "status": "completed",
                "similarity_analysis": similarity_analysis,
                "ai_analysis": ai_analysis,
                "plagiarism_analysis": plagiarism_analysis,
                "report_html": html_report,
                "processing_info": {
            "filename": filename,
            "content_length": len(content),
            "word_count": len(content.split()),
            "saved_txt_path": saved_txt_path,
            "corpus_dir": CORPUS_DIR,
            "corpus_files_compared": similarity_analysis['total_comparisons']
                },
                "errors": []
            }
            
        except Exception as e:
        logger.error(f"Error in analysis: {e}")
            return {
                "status": "failed",
                "similarity_analysis": {"similarity_score": 0, "total_comparisons": 0, "detailed_results": []},
                "ai_analysis": {"ai_percentage": 0, "interpretation": "Analysis failed", "chunks_analyzed": 0, "confidence": 0},
                "plagiarism_analysis": {"plagiarism_score": 0, "risk_level": "unknown", "contributing_factors": ["analysis_error"]},
                "report_html": f"<html><body><h1>Analysis Failed</h1><p>Error: {str(e)}</p></body></html>",
                "processing_info": {"filename": "unknown", "content_length": 0, "word_count": 0},
                "errors": [str(e)]
            }


# Initialize API instance
api = PlagiarismAPI()

@app.route('/grade/analyze', methods=['POST'])
def analyze_submission():
    """Main API endpoint for plagiarism analysis with PDF/OCR processing"""
    try:
        # Validate request
        if not request.is_json:
            return jsonify({"error": "Content-Type must be application/json"}), 400
            
        payload = request.get_json()
        
        # Validate required fields
        required_fields = ['student_id', 'student_name', 'file_url', 'assignment_id', 'classroom_name']
        for field in required_fields:
            if field not in payload:
                return jsonify({"error": f"Missing required field: {field}"}), 400
        
        # Perform analysis with integrated PDF/OCR processing
    result = api.analyze_document(payload)
        
        return jsonify(result)
        
    except Exception as e:
    logger.error(f"API error: {e}")
        return jsonify({
            "status": "failed",
            "similarity_analysis": {"similarity_score": 0, "total_comparisons": 0, "detailed_results": []},
            "ai_analysis": {"ai_percentage": 0, "interpretation": "Analysis failed", "chunks_analyzed": 0, "confidence": 0},
            "plagiarism_analysis": {"plagiarism_score": 0, "risk_level": "unknown", "contributing_factors": ["api_error"]},
            "report_html": f"<html><body><h1>Analysis Failed</h1><p>API Error: {str(e)}</p></body></html>",
            "processing_info": {"filename": "unknown", "content_length": 0, "word_count": 0},
            "errors": [str(e)]
        }), 500


@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "service": "Enhanced Plagiarism Detection API with PDF/OCR Processing",
        "capabilities": [
            "PDF text extraction",
            "OCR for image-based documents",
            "DOCX/PPTX processing",
            "Image preprocessing",
            "Multi-format support"
        ]
    })


@app.route('/test-processing', methods=['POST'])
def test_file_processing():
    """Test endpoint for file processing capabilities"""
    try:
        if not request.is_json:
            return jsonify({"error": "Content-Type must be application/json"}), 400
            
        payload = request.get_json()
        
        if 'file_url' not in payload:
            return jsonify({"error": "Missing required field: file_url"}), 400
        
    # Test file processing
    logger.info(f"Test-processing request for URL: {payload.get('file_url')}")
    content, filename = api.download_and_process_file(payload['file_url'])
        
        return jsonify({
            "status": "success",
            "filename": filename,
            "content_preview": content[:500] + "..." if len(content) > 500 else content,
            "content_length": len(content),
            "word_count": len(content.split())
        })
        
    except Exception as e:
    logger.error(f"Test processing error: {e}")
        return jsonify({
            "status": "failed",
            "error": str(e)
        }), 500


@app.route('/user/<student_id>/scores', methods=['GET'])
def get_user_scores(student_id):
    """Get historical scores for a user"""
    global user_scores_log
    
    if student_id in user_scores_log:
        return jsonify({
            "student_id": student_id,
            "scores_history": user_scores_log[student_id]
        })
    else:
        return jsonify({
            "student_id": student_id,
            "scores_history": [],
            "message": "No scores found for this user"
        })


@app.route('/stats', methods=['GET'])
def get_stats():
    """Get system statistics"""
    global user_scores_log
    
    total_users = len(user_scores_log)
    total_analyses = sum(len(scores) for scores in user_scores_log.values())
    try:
        corpus_count = len([n for n in os.listdir(CORPUS_DIR) if n.lower().endswith('.txt')])
    except Exception:
        corpus_count = 0
    
    return jsonify({
        "total_users_analyzed": total_users,
        "total_analyses_performed": total_analyses,
        "corpus_text_files": corpus_count,
        "system_status": "operational",
        "features": [
            "PDF text extraction with PyMuPDF",
            "OCR processing with Tesseract",
            "Image preprocessing with OpenCV",
            "DOCX/PPTX support",
            "Multi-format document processing",
            "Persistent text corpus and folder-based similarity"
        ]
    })


if __name__ == '__main__':
    # Create logs & corpus directory
    os.makedirs(LOGS_DIR, exist_ok=True)
    os.makedirs(CORPUS_DIR, exist_ok=True)
    
    print("Starting Enhanced Plagiarism Detection API Service with PDF/OCR Processing...")
    print("Server will run on http://localhost:5000")
    print("\nFeatures enabled:")
    print("✓ PDF text extraction with PyMuPDF")
    print("✓ OCR processing with Tesseract")
    print("✓ Image preprocessing with OpenCV/PIL")
    print("✓ DOCX/PPTX document support")
    print("✓ Multi-format file processing")
    print("\nAvailable endpoints:")
    print("POST /grade/analyze - Main plagiarism analysis endpoint")
    print("POST /test-processing - Test file processing capabilities")
    print("GET /health - Health check with capabilities")
    print("GET /user/<student_id>/scores - Get user's score history")
    print("GET /stats - Get system statistics")
    
    # Test Tesseract installation on startup
    try:
        pytesseract.get_tesseract_version()
        print(f"✓ Tesseract OCR ready: {pytesseract.get_tesseract_version()}")
    except Exception as e:
        print(f"⚠ Tesseract warning: {e}")
        print("  OCR functionality may be limited")
    
    app.run(host='0.0.0.0', port=5000, debug=True)